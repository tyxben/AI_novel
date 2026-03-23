"""Tests for HTMLToPPTXConverter — HTML slide → PPTX conversion.

All tests are SYNC and mock Playwright entirely (no real browser launched).
BeautifulSoup is mocked or used directly depending on the test.
"""

from __future__ import annotations

import textwrap
from pathlib import Path
from unittest.mock import MagicMock, call, patch

import pytest
from pptx import Presentation
from pptx.util import Inches

# ---------------------------------------------------------------------------
# Sample HTML for tests
# ---------------------------------------------------------------------------

SAMPLE_HTML = textwrap.dedent("""\
    <!DOCTYPE html>
    <html lang="zh-CN">
    <head><meta charset="UTF-8"><title>Test</title></head>
    <body>
    <section class="slide" data-page="1">
        <h1>标题页</h1>
        <p class="subtitle">副标题文字</p>
    </section>
    <section class="slide" data-page="2">
        <h1>第一章</h1>
        <ul>
            <li>要点 A</li>
            <li>要点 B</li>
        </ul>
        <p>正文段落内容</p>
    </section>
    <section class="slide" data-page="3">
        <h2>小节标题</h2>
        <p>另一段内容</p>
    </section>
    </body>
    </html>
""")

SAMPLE_HTML_NO_SLIDES = textwrap.dedent("""\
    <!DOCTYPE html>
    <html><body><div>No slides here</div></body></html>
""")

EMPTY_SLIDE_HTML = textwrap.dedent("""\
    <!DOCTYPE html>
    <html><body>
    <section class="slide" data-page="1">
    </section>
    </body></html>
""")


# ---------------------------------------------------------------------------
# Helper: create a minimal PNG file (1x1 red pixel)
# ---------------------------------------------------------------------------

def _create_test_png(path: Path) -> Path:
    """Create a minimal valid PNG file at *path*."""
    try:
        from PIL import Image

        img = Image.new("RGB", (1280, 720), color="red")
        img.save(str(path), format="PNG")
    except ImportError:
        # Fallback: write a minimal 1x1 PNG manually (8-bit RGB, no filter)
        import struct
        import zlib

        def _chunk(chunk_type: bytes, data: bytes) -> bytes:
            raw = chunk_type + data
            return struct.pack(">I", len(data)) + raw + struct.pack(">I", zlib.crc32(raw) & 0xFFFFFFFF)

        signature = b"\x89PNG\r\n\x1a\n"
        ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
        raw_data = zlib.compress(b"\x00\xff\x00\x00")
        png_bytes = signature + _chunk(b"IHDR", ihdr) + _chunk(b"IDAT", raw_data) + _chunk(b"IEND", b"")
        path.write_bytes(png_bytes)
    return path


# ---------------------------------------------------------------------------
# Tests: Initialization
# ---------------------------------------------------------------------------


class TestInit:
    """Test HTMLToPPTXConverter initialization."""

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", False)
    def test_init_without_playwright(self):
        """Raise ImportError with helpful message when Playwright is missing."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        with pytest.raises(ImportError, match="Playwright is required"):
            HTMLToPPTXConverter()

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_init_with_playwright(self):
        """Successful instantiation when Playwright is available."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(workspace="/tmp/test_ws")
        assert converter.workspace == Path("/tmp/test_ws")

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    @patch("src.ppt.html_to_pptx._HAS_BS4", False)
    def test_init_extract_text_disabled_without_bs4(self):
        """Text extraction silently disabled when BS4 is missing."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(extract_text=True)
        assert converter.extract_text is False

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    @patch("src.ppt.html_to_pptx._HAS_BS4", True)
    def test_init_extract_text_enabled(self):
        """Text extraction enabled when BS4 is available and requested."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(extract_text=True)
        assert converter.extract_text is True

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    @patch("src.ppt.html_to_pptx._HAS_BS4", True)
    def test_init_extract_text_explicitly_disabled(self):
        """Text extraction disabled when caller sets extract_text=False."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(extract_text=False)
        assert converter.extract_text is False


# ---------------------------------------------------------------------------
# Tests: Text extraction
# ---------------------------------------------------------------------------


class TestExtractTextLayers:
    """Test _extract_text_layers with real BeautifulSoup."""

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_extract_text_layers(self):
        """Extract structured text from multi-slide HTML."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter()

        # If BS4 is not installed, skip this test
        if not converter.extract_text:
            pytest.skip("BeautifulSoup not installed")

        layers = converter._extract_text_layers(SAMPLE_HTML)

        assert len(layers) == 3

        # Slide 1: title + subtitle
        assert "# 标题页" in layers[0]
        assert "## 副标题文字" in layers[0]

        # Slide 2: title + bullets + paragraph
        assert "# 第一章" in layers[1]
        assert "- 要点 A" in layers[1]
        assert "- 要点 B" in layers[1]
        assert "正文段落内容" in layers[1]

        # Slide 3: h2 heading + paragraph
        assert "## 小节标题" in layers[2]
        assert "另一段内容" in layers[2]

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_extract_text_empty_slide(self):
        """Empty slide returns empty string."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter()
        if not converter.extract_text:
            pytest.skip("BeautifulSoup not installed")

        layers = converter._extract_text_layers(EMPTY_SLIDE_HTML)
        assert len(layers) == 1
        assert layers[0] == ""

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    @patch("src.ppt.html_to_pptx._HAS_BS4", False)
    def test_extract_text_no_bs4(self):
        """Returns empty list when BS4 is not available."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter()
        layers = converter._extract_text_layers(SAMPLE_HTML)
        assert layers == []


# ---------------------------------------------------------------------------
# Tests: PPTX assembly
# ---------------------------------------------------------------------------


class TestAssemblePptx:
    """Test _assemble_pptx with real python-pptx."""

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_assemble_pptx(self, tmp_path: Path):
        """Create PPTX from mock PNG files — verify slide count and dimensions."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(workspace=str(tmp_path))

        # Create test PNG files
        img_dir = tmp_path / "imgs"
        img_dir.mkdir()
        img1 = _create_test_png(img_dir / "slide_001.png")
        img2 = _create_test_png(img_dir / "slide_002.png")

        output = tmp_path / "output.pptx"
        result = converter._assemble_pptx([img1, img2], [], output, None)

        assert result == output
        assert output.exists()

        # Verify PPTX content
        prs = Presentation(str(output))
        assert len(prs.slides) == 2
        assert prs.slide_width == Inches(13.333)
        assert prs.slide_height == Inches(7.5)

        # Each slide should have exactly one picture shape
        for slide in prs.slides:
            pics = [s for s in slide.shapes if s.shape_type is not None]
            assert len(pics) >= 1  # At least the image

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_assemble_pptx_with_notes(self, tmp_path: Path):
        """Verify notes text is added to slides."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(workspace=str(tmp_path))

        img_dir = tmp_path / "imgs"
        img_dir.mkdir()
        img1 = _create_test_png(img_dir / "slide_001.png")
        img2 = _create_test_png(img_dir / "slide_002.png")

        notes = ["# 标题页\n## 副标题", "- 要点 A\n- 要点 B"]
        output = tmp_path / "output_notes.pptx"
        converter._assemble_pptx([img1, img2], notes, output, None)

        prs = Presentation(str(output))
        slide1_notes = prs.slides[0].notes_slide.notes_text_frame.text
        slide2_notes = prs.slides[1].notes_slide.notes_text_frame.text

        assert "标题页" in slide1_notes
        assert "副标题" in slide1_notes
        assert "要点 A" in slide2_notes
        assert "要点 B" in slide2_notes

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_assemble_pptx_empty_screenshots(self, tmp_path: Path):
        """Empty screenshot list creates an empty PPTX (no slides)."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(workspace=str(tmp_path))
        output = tmp_path / "empty.pptx"
        converter._assemble_pptx([], [], output, None)

        prs = Presentation(str(output))
        assert len(prs.slides) == 0

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_assemble_pptx_fewer_notes_than_slides(self, tmp_path: Path):
        """When text_layers is shorter than screenshot_paths, extra slides get no notes."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(workspace=str(tmp_path))

        img_dir = tmp_path / "imgs"
        img_dir.mkdir()
        img1 = _create_test_png(img_dir / "slide_001.png")
        img2 = _create_test_png(img_dir / "slide_002.png")
        img3 = _create_test_png(img_dir / "slide_003.png")

        notes = ["# 标题"]  # Only one note for three slides
        output = tmp_path / "partial_notes.pptx"
        converter._assemble_pptx([img1, img2, img3], notes, output, None)

        prs = Presentation(str(output))
        assert len(prs.slides) == 3

        # First slide has notes
        assert "标题" in prs.slides[0].notes_slide.notes_text_frame.text

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_assemble_creates_output_directory(self, tmp_path: Path):
        """Output directory is created automatically if it does not exist."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(workspace=str(tmp_path))

        img_dir = tmp_path / "imgs"
        img_dir.mkdir()
        img1 = _create_test_png(img_dir / "slide_001.png")

        # Output in a non-existent subdirectory
        output = tmp_path / "nested" / "dir" / "result.pptx"
        converter._assemble_pptx([img1], [], output, None)
        assert output.exists()


# ---------------------------------------------------------------------------
# Tests: Cleanup
# ---------------------------------------------------------------------------


class TestCleanup:
    """Test _cleanup_screenshots."""

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_cleanup_screenshots(self, tmp_path: Path):
        """Verify temp files are removed."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(workspace=str(tmp_path))

        screenshots_dir = tmp_path / "screenshots"
        screenshots_dir.mkdir()
        (screenshots_dir / "slide_001.png").write_bytes(b"\x89PNG")
        (screenshots_dir / "slide_002.png").write_bytes(b"\x89PNG")

        assert screenshots_dir.exists()
        converter._cleanup_screenshots(screenshots_dir)
        assert not screenshots_dir.exists()

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_cleanup_nonexistent_dir(self, tmp_path: Path):
        """Cleaning up a non-existent directory is a no-op (no error)."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(workspace=str(tmp_path))
        converter._cleanup_screenshots(tmp_path / "does_not_exist")
        # Should not raise


# ---------------------------------------------------------------------------
# Tests: Slide counting
# ---------------------------------------------------------------------------


class TestCountSlides:
    """Test _count_slides regex-based counting."""

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_count_slides_normal(self):
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter()
        assert converter._count_slides(SAMPLE_HTML) == 3

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_count_slides_zero(self):
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter()
        assert converter._count_slides(SAMPLE_HTML_NO_SLIDES) == 0

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_count_slides_empty_html(self):
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter()
        assert converter._count_slides("") == 0

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_count_slides_single(self):
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter()
        assert converter._count_slides(EMPTY_SLIDE_HTML) == 1


# ---------------------------------------------------------------------------
# Tests: Full convert flow (Playwright mocked)
# ---------------------------------------------------------------------------


class TestConvert:
    """Test the full convert() method with Playwright fully mocked."""

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_convert_no_slides_raises(self, tmp_path: Path):
        """HTML with no .slide sections raises ValueError."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(workspace=str(tmp_path))
        html_file = tmp_path / "empty.html"
        html_file.write_text(SAMPLE_HTML_NO_SLIDES, encoding="utf-8")

        with pytest.raises(ValueError, match="No slides found"):
            converter.convert(html_file, tmp_path / "out.pptx")

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_convert_file_not_found(self, tmp_path: Path):
        """Raises FileNotFoundError for missing HTML file."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(workspace=str(tmp_path))
        with pytest.raises(FileNotFoundError, match="HTML file not found"):
            converter.convert(tmp_path / "nonexistent.html", tmp_path / "out.pptx")

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_convert_full_flow(self, tmp_path: Path):
        """Full convert flow with mocked Playwright — verify PPTX is created."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(workspace=str(tmp_path))

        # Write sample HTML
        html_file = tmp_path / "slides.html"
        html_file.write_text(SAMPLE_HTML, encoding="utf-8")
        output_path = tmp_path / "result.pptx"

        # Create mock PNGs that _capture_screenshots will "produce"
        def mock_capture(html_path, output_dir, total_pages, callback):
            paths = []
            for i in range(total_pages):
                p = output_dir / f"slide_{i+1:03d}.png"
                _create_test_png(p)
                paths.append(p)
            return paths

        with patch.object(converter, "_capture_screenshots", side_effect=mock_capture):
            result = converter.convert(html_file, output_path)

        assert result == str(output_path)
        assert output_path.exists()

        prs = Presentation(str(output_path))
        assert len(prs.slides) == 3

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_progress_callback_called(self, tmp_path: Path):
        """Verify progress callback is called with correct arguments."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(workspace=str(tmp_path))

        html_file = tmp_path / "slides.html"
        html_file.write_text(SAMPLE_HTML, encoding="utf-8")
        output_path = tmp_path / "result.pptx"

        callback = MagicMock()

        def mock_capture(html_path, output_dir, total_pages, cb):
            paths = []
            for i in range(total_pages):
                p = output_dir / f"slide_{i+1:03d}.png"
                _create_test_png(p)
                if cb:
                    cb(i + 1, total_pages, f"截图第 {i+1}/{total_pages} 页")
                paths.append(p)
            return paths

        with patch.object(converter, "_capture_screenshots", side_effect=mock_capture):
            converter.convert(html_file, output_path, progress_callback=callback)

        # Callback should be called for screenshots (3) + assembly (3) = 6 times
        assert callback.call_count >= 3

        # Check that all screenshot calls were made
        screenshot_calls = [
            c for c in callback.call_args_list if "截图" in str(c)
        ]
        assert len(screenshot_calls) == 3

        # Check that assembly calls were made (page_num, total, message)
        assembly_calls = [
            c for c in callback.call_args_list if "组装" in str(c)
        ]
        assert len(assembly_calls) == 3

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_convert_cleanup_on_success(self, tmp_path: Path):
        """Verify temp screenshots are cleaned up after successful conversion."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(workspace=str(tmp_path))

        html_file = tmp_path / "slides.html"
        html_file.write_text(SAMPLE_HTML, encoding="utf-8")
        output_path = tmp_path / "result.pptx"

        created_dirs: list[Path] = []

        def mock_capture(html_path, output_dir, total_pages, callback):
            created_dirs.append(output_dir)
            paths = []
            for i in range(total_pages):
                p = output_dir / f"slide_{i+1:03d}.png"
                _create_test_png(p)
                paths.append(p)
            return paths

        with patch.object(converter, "_capture_screenshots", side_effect=mock_capture):
            converter.convert(html_file, output_path)

        # The temp directory should have been cleaned up
        assert len(created_dirs) == 1
        assert not created_dirs[0].exists()

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_convert_cleanup_on_error(self, tmp_path: Path):
        """Verify temp screenshots are cleaned up even if assembly fails."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(workspace=str(tmp_path))

        html_file = tmp_path / "slides.html"
        html_file.write_text(SAMPLE_HTML, encoding="utf-8")
        output_path = tmp_path / "result.pptx"

        created_dirs: list[Path] = []

        def mock_capture(html_path, output_dir, total_pages, callback):
            created_dirs.append(output_dir)
            paths = []
            for i in range(total_pages):
                p = output_dir / f"slide_{i+1:03d}.png"
                _create_test_png(p)
                paths.append(p)
            return paths

        with patch.object(converter, "_capture_screenshots", side_effect=mock_capture):
            with patch.object(converter, "_assemble_pptx", side_effect=RuntimeError("boom")):
                with pytest.raises(RuntimeError, match="boom"):
                    converter.convert(html_file, output_path)

        # Cleanup should still have happened
        assert len(created_dirs) == 1
        assert not created_dirs[0].exists()


# ---------------------------------------------------------------------------
# Tests: Capture screenshots (Playwright fully mocked)
# ---------------------------------------------------------------------------


class TestCaptureScreenshots:
    """Test _capture_screenshots with Playwright fully mocked."""

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_capture_screenshots_calls_playwright(self, tmp_path: Path):
        """Verify Playwright interactions: launch, goto, evaluate, screenshot."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(workspace=str(tmp_path))

        html_file = tmp_path / "test.html"
        html_file.write_text(SAMPLE_HTML, encoding="utf-8")
        output_dir = tmp_path / "screenshots"
        output_dir.mkdir()

        # Build mock chain: sync_playwright → pw → browser → page → element
        mock_element = MagicMock()
        mock_page = MagicMock()
        mock_page.query_selector.return_value = mock_element
        mock_browser = MagicMock()
        mock_browser.new_page.return_value = mock_page
        mock_pw = MagicMock()
        mock_pw.chromium.launch.return_value = mock_browser

        mock_context_manager = MagicMock()
        mock_context_manager.__enter__ = MagicMock(return_value=mock_pw)
        mock_context_manager.__exit__ = MagicMock(return_value=False)

        with patch("playwright.sync_api.sync_playwright", return_value=mock_context_manager):
            paths = converter._capture_screenshots(html_file, output_dir, 3, None)

        # Verify interactions
        mock_pw.chromium.launch.assert_called_once_with(headless=True)
        mock_browser.new_page.assert_called_once()
        mock_page.goto.assert_called_once()

        # showSlide called for each page (0-based)
        evaluate_calls = mock_page.evaluate.call_args_list
        assert len(evaluate_calls) == 3
        assert evaluate_calls[0] == call("showSlide(0)")
        assert evaluate_calls[1] == call("showSlide(1)")
        assert evaluate_calls[2] == call("showSlide(2)")

        # wait_for_timeout called 3 times
        assert mock_page.wait_for_timeout.call_count == 3

        # screenshot called 3 times on the element
        assert mock_element.screenshot.call_count == 3

        # Browser closed
        mock_browser.close.assert_called_once()

        # Returned paths
        assert len(paths) == 3
        for p in paths:
            assert p.suffix == ".png"

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_capture_screenshots_fallback_no_element(self, tmp_path: Path):
        """Falls back to page.screenshot when slide element is not found."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(workspace=str(tmp_path))

        html_file = tmp_path / "test.html"
        html_file.write_text(EMPTY_SLIDE_HTML, encoding="utf-8")
        output_dir = tmp_path / "screenshots"
        output_dir.mkdir()

        mock_page = MagicMock()
        mock_page.query_selector.return_value = None
        mock_page.query_selector_all.return_value = []

        mock_browser = MagicMock()
        mock_browser.new_page.return_value = mock_page
        mock_pw = MagicMock()
        mock_pw.chromium.launch.return_value = mock_browser

        mock_context_manager = MagicMock()
        mock_context_manager.__enter__ = MagicMock(return_value=mock_pw)
        mock_context_manager.__exit__ = MagicMock(return_value=False)

        with patch("playwright.sync_api.sync_playwright", return_value=mock_context_manager):
            paths = converter._capture_screenshots(html_file, output_dir, 1, None)

        # Should fall back to full-page screenshot
        mock_page.screenshot.assert_called_once()
        assert len(paths) == 1

    @patch("src.ppt.html_to_pptx._HAS_PLAYWRIGHT", True)
    def test_capture_screenshots_with_callback(self, tmp_path: Path):
        """Verify progress callback is invoked during capture."""
        from src.ppt.html_to_pptx import HTMLToPPTXConverter

        converter = HTMLToPPTXConverter(workspace=str(tmp_path))

        html_file = tmp_path / "test.html"
        html_file.write_text(SAMPLE_HTML, encoding="utf-8")
        output_dir = tmp_path / "screenshots"
        output_dir.mkdir()

        mock_element = MagicMock()
        mock_page = MagicMock()
        mock_page.query_selector.return_value = mock_element
        mock_browser = MagicMock()
        mock_browser.new_page.return_value = mock_page
        mock_pw = MagicMock()
        mock_pw.chromium.launch.return_value = mock_browser

        mock_context_manager = MagicMock()
        mock_context_manager.__enter__ = MagicMock(return_value=mock_pw)
        mock_context_manager.__exit__ = MagicMock(return_value=False)

        callback = MagicMock()

        with patch("playwright.sync_api.sync_playwright", return_value=mock_context_manager):
            converter._capture_screenshots(html_file, output_dir, 3, callback)

        assert callback.call_count == 3
        callback.assert_any_call(1, 3, "截图第 1/3 页")
        callback.assert_any_call(2, 3, "截图第 2/3 页")
        callback.assert_any_call(3, 3, "截图第 3/3 页")
