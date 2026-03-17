"""Tests for PPTRenderer — covers every layout, utility methods, and save/load."""

from __future__ import annotations

import os
import tempfile

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from src.ppt.models import (
    ColorScheme,
    DecorationSpec,
    FontSpec,
    LayoutType,
    SlideContent,
    SlideDesign,
    SlideSpec,
    ThemeConfig,
)
from src.ppt.ppt_renderer import PPTRenderer


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


def _make_color_scheme(**overrides) -> ColorScheme:
    defaults = dict(
        primary="#1A237E",
        secondary="#00BCD4",
        accent="#FDD835",
        text="#424242",
        background="#FFFFFF",
    )
    defaults.update(overrides)
    return ColorScheme(**defaults)


def _make_font(size: int = 20, bold: bool = False, color: str = "#424242") -> FontSpec:
    return FontSpec(size=size, bold=bold, color=color, family="Arial")


def _make_theme(**overrides) -> ThemeConfig:
    defaults = dict(
        name="test_modern",
        display_name="Test Modern",
        colors=_make_color_scheme(),
        title_font=_make_font(32, True, "#1A237E"),
        body_font=_make_font(18, False, "#424242"),
        note_font=_make_font(12, False, "#757575"),
    )
    defaults.update(overrides)
    return ThemeConfig(**defaults)


def _make_design(layout: LayoutType, **overrides) -> SlideDesign:
    defaults = dict(
        layout=layout,
        colors=_make_color_scheme(),
        title_font=_make_font(32, True, "#1A237E"),
        body_font=_make_font(18, False, "#424242"),
        note_font=_make_font(12, False, "#757575"),
        decoration=DecorationSpec(),
    )
    defaults.update(overrides)
    return SlideDesign(**defaults)


def _make_slide(
    layout: LayoutType,
    page: int = 1,
    title: str = "Test Title",
    subtitle: str | None = None,
    bullets: list[str] | None = None,
    body_text: str | None = None,
    data_value: str | None = None,
    data_label: str | None = None,
    speaker_notes: str = "",
    image_path: str | None = None,
) -> SlideSpec:
    return SlideSpec(
        page_number=page,
        content=SlideContent(
            title=title,
            subtitle=subtitle,
            bullet_points=bullets or [],
            body_text=body_text,
            speaker_notes=speaker_notes,
            data_value=data_value,
            data_label=data_label,
        ),
        design=_make_design(layout),
        image_path=image_path,
    )


@pytest.fixture
def theme() -> ThemeConfig:
    return _make_theme()


@pytest.fixture
def renderer(theme) -> PPTRenderer:
    return PPTRenderer(theme)


# ---------------------------------------------------------------------------
# Helper: save, reload and return Presentation
# ---------------------------------------------------------------------------


def _render_and_reload(renderer: PPTRenderer, slides: list[SlideSpec]) -> Presentation:
    """Render slides, save to temp file, reload and return Presentation."""
    renderer.render(slides)
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
    try:
        renderer.save(path)
        return Presentation(path)
    finally:
        os.unlink(path)


# ===========================================================================
# Test: PPTRenderer initialisation
# ===========================================================================


class TestRendererInit:
    def test_slide_dimensions(self, renderer):
        assert renderer.prs.slide_width == Inches(13.333)
        assert renderer.prs.slide_height == Inches(7.5)

    def test_theme_stored(self, renderer, theme):
        assert renderer.theme is theme


# ===========================================================================
# Test: Utility methods
# ===========================================================================


class TestUtilities:
    def test_hex_to_rgb_standard(self):
        assert PPTRenderer._hex_to_rgb("#FF0000") == RGBColor(255, 0, 0)
        assert PPTRenderer._hex_to_rgb("#00ff00") == RGBColor(0, 255, 0)

    def test_hex_to_rgb_no_hash(self):
        assert PPTRenderer._hex_to_rgb("0000FF") == RGBColor(0, 0, 255)

    def test_hex_to_rgb_invalid_fallback(self):
        # Invalid hex should fall back gracefully
        result = PPTRenderer._hex_to_rgb("ZZZ")
        assert isinstance(result, RGBColor)

    def test_lighten_or_darken(self):
        darker = PPTRenderer._lighten_or_darken("#808080", darken=16)
        assert darker == "#707070"

        lighter = PPTRenderer._lighten_or_darken("#808080", lighten=16)
        assert lighter == "#909090"

    def test_lighten_clamps(self):
        result = PPTRenderer._lighten_or_darken("#FFFFFF", lighten=100)
        assert result == "#FFFFFF"

    def test_darken_clamps(self):
        result = PPTRenderer._lighten_or_darken("#000000", darken=100)
        assert result == "#000000"

    def test_muted_color_full_opacity(self):
        assert PPTRenderer._muted_color("#FF0000", 1.0) == "#FF0000"

    def test_muted_color_zero_opacity(self):
        assert PPTRenderer._muted_color("#FF0000", 0.0) == "#FFFFFF"

    def test_muted_color_half(self):
        result = PPTRenderer._muted_color("#000000", 0.5)
        # Should be roughly #7F7F7F
        h = result.lstrip("#")
        r = int(h[0:2], 16)
        assert 126 <= r <= 128


class TestAddTextBox:
    def test_basic_text_box(self, renderer):
        slide = renderer._add_blank_slide()
        shape = renderer._add_text_box(
            slide, Inches(1), Inches(1), Inches(4), Inches(1),
            "Hello World", font_size=24, bold=True, color="#FF0000",
        )
        assert shape.text_frame.paragraphs[0].runs[0].text == "Hello World"
        assert shape.text_frame.paragraphs[0].runs[0].font.bold is True
        assert shape.text_frame.paragraphs[0].runs[0].font.size.pt == 24

    def test_italic_text(self, renderer):
        slide = renderer._add_blank_slide()
        shape = renderer._add_text_box(
            slide, Inches(1), Inches(1), Inches(4), Inches(1),
            "Italic text", italic=True,
        )
        assert shape.text_frame.paragraphs[0].runs[0].font.italic is True


class TestAddShape:
    def test_filled_rectangle(self, renderer):
        slide = renderer._add_blank_slide()
        shape = renderer._add_shape(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(1), Inches(3), Inches(2),
            fill_color="#FF0000",
        )
        assert shape.fill.fore_color.rgb == RGBColor(255, 0, 0)

    def test_shape_no_fill(self, renderer):
        slide = renderer._add_blank_slide()
        shape = renderer._add_shape(
            slide, MSO_SHAPE.OVAL,
            Inches(1), Inches(1), Inches(1), Inches(1),
        )
        # Should not raise
        assert shape is not None

    def test_shape_with_line(self, renderer):
        slide = renderer._add_blank_slide()
        shape = renderer._add_shape(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(1), Inches(1), Inches(2), Inches(2),
            line_color="#00FF00", line_width=2.0,
        )
        assert shape.line.color.rgb == RGBColor(0, 255, 0)


class TestAccentLine:
    def test_adds_thin_rectangle(self, renderer):
        slide = renderer._add_blank_slide()
        before = len(slide.shapes)
        renderer._add_accent_line(slide, Inches(1), Inches(1), Inches(3))
        assert len(slide.shapes) == before + 1


class TestOverlay:
    def test_overlay_shape_added(self, renderer):
        slide = renderer._add_blank_slide()
        before = len(slide.shapes)
        renderer._add_overlay(slide, color="#000000", opacity=0.5)
        assert len(slide.shapes) == before + 1


class TestGradientRect:
    def test_gradient_shape_added(self, renderer):
        slide = renderer._add_blank_slide()
        before = len(slide.shapes)
        renderer._add_gradient_rect(
            slide, Inches(0), Inches(0), Inches(5), Inches(2),
            "#1A237E", "#FDD835",
        )
        assert len(slide.shapes) == before + 1


class TestSetBackground:
    def test_sets_background(self, renderer):
        slide = renderer._add_blank_slide()
        renderer._set_slide_background(slide, "#F5F5F5")
        # Should not raise; background fill is set
        assert slide.background.fill.fore_color.rgb == RGBColor(0xF5, 0xF5, 0xF5)


class TestFooter:
    def test_footer_with_text_and_page(self, renderer):
        slide = renderer._add_blank_slide()
        before = len(slide.shapes)
        renderer._add_footer(slide, text="Confidential", page_number=3)
        assert len(slide.shapes) == before + 1

    def test_footer_empty_no_shape(self, renderer):
        slide = renderer._add_blank_slide()
        before = len(slide.shapes)
        renderer._add_footer(slide, text="", page_number=0)
        assert len(slide.shapes) == before  # nothing added


# ===========================================================================
# Test: Each layout renders without error and produces correct slide count
# ===========================================================================


class TestLayoutRendering:
    """Each layout should produce exactly one slide and be saveable."""

    def _assert_one_slide(self, renderer, slide_spec):
        prs = _render_and_reload(renderer, [slide_spec])
        assert len(prs.slides) == 1

    def test_title_hero(self, renderer):
        self._assert_one_slide(renderer, _make_slide(
            LayoutType.TITLE_HERO,
            title="2024 Product Report",
            subtitle="Data-Driven Insights",
            body_text="March 2024 | Marketing Team",
        ))

    def test_section_divider(self, renderer):
        self._assert_one_slide(renderer, _make_slide(
            LayoutType.SECTION_DIVIDER,
            title="Market Overview",
        ))

    def test_text_left_image_right(self, renderer):
        self._assert_one_slide(renderer, _make_slide(
            LayoutType.TEXT_LEFT_IMAGE_RIGHT,
            title="Key Findings",
            bullets=["Point A", "Point B", "Point C"],
        ))

    def test_image_left_text_right(self, renderer):
        self._assert_one_slide(renderer, _make_slide(
            LayoutType.IMAGE_LEFT_TEXT_RIGHT,
            title="Visual Analysis",
            bullets=["Finding 1", "Finding 2"],
        ))

    def test_full_image_overlay(self, renderer):
        self._assert_one_slide(renderer, _make_slide(
            LayoutType.FULL_IMAGE_OVERLAY,
            title="Bold Statement",
            body_text="Supporting text here.",
        ))

    def test_three_columns(self, renderer):
        self._assert_one_slide(renderer, _make_slide(
            LayoutType.THREE_COLUMNS,
            title="Three Pillars",
            bullets=["Speed: Blazing fast", "Quality: Top notch", "Cost: Affordable"],
        ))

    def test_quote_page(self, renderer):
        self._assert_one_slide(renderer, _make_slide(
            LayoutType.QUOTE_PAGE,
            title="Inspiration",
            subtitle="Steve Jobs",
            body_text="Stay hungry, stay foolish.",
        ))

    def test_data_highlight(self, renderer):
        self._assert_one_slide(renderer, _make_slide(
            LayoutType.DATA_HIGHLIGHT,
            title="Growth Metrics",
            data_value="30%",
            data_label="Year-over-Year Growth",
            bullets=["Exceeded expectations by 5 percentage points"],
        ))

    def test_timeline(self, renderer):
        self._assert_one_slide(renderer, _make_slide(
            LayoutType.TIMELINE,
            title="Company History",
            bullets=[
                "2020: Founded",
                "2021: Series A",
                "2022: 100 employees",
                "2023: IPO",
            ],
        ))

    def test_timeline_empty(self, renderer):
        """Timeline with no bullets should still render."""
        self._assert_one_slide(renderer, _make_slide(
            LayoutType.TIMELINE,
            title="Empty Timeline",
            bullets=[],
        ))

    def test_bullet_with_icons(self, renderer):
        self._assert_one_slide(renderer, _make_slide(
            LayoutType.BULLET_WITH_ICONS,
            title="Key Features",
            bullets=[
                "Performance: 10x faster than before",
                "Security: Enterprise-grade encryption",
                "Scalability: Handles millions of requests",
            ],
        ))

    def test_comparison(self, renderer):
        self._assert_one_slide(renderer, _make_slide(
            LayoutType.COMPARISON,
            title="Options Comparison",
            subtitle="Plan A vs Plan B",
            bullets=["Cheaper", "Simpler", "Faster", "More powerful", "Better support", "Enterprise-ready"],
        ))

    def test_closing(self, renderer):
        self._assert_one_slide(renderer, _make_slide(
            LayoutType.CLOSING,
            title="Thank You!",
            subtitle="Questions?",
            body_text="email@example.com | @handle",
        ))


# ===========================================================================
# Test: Multi-slide rendering
# ===========================================================================


class TestMultiSlideRendering:
    def test_renders_multiple_slides(self, renderer):
        slides = [
            _make_slide(LayoutType.TITLE_HERO, page=1, title="Cover"),
            _make_slide(LayoutType.SECTION_DIVIDER, page=2, title="Section 1"),
            _make_slide(LayoutType.BULLET_WITH_ICONS, page=3, title="Points",
                        bullets=["A", "B", "C"]),
            _make_slide(LayoutType.DATA_HIGHLIGHT, page=4, title="Data",
                        data_value="99%", data_label="Accuracy"),
            _make_slide(LayoutType.CLOSING, page=5, title="Bye"),
        ]
        prs = _render_and_reload(renderer, slides)
        assert len(prs.slides) == 5

    def test_speaker_notes_preserved(self, renderer):
        slide = _make_slide(
            LayoutType.SECTION_DIVIDER,
            title="Notes Test",
            speaker_notes="Remember to mention the quarterly results.",
        )
        prs = _render_and_reload(renderer, [slide])
        notes_text = prs.slides[0].notes_slide.notes_text_frame.text
        assert "quarterly results" in notes_text


# ===========================================================================
# Test: Footer behaviour
# ===========================================================================


class TestFooterBehaviour:
    def test_title_hero_no_footer(self, renderer):
        """Title page should not have a footer."""
        slide_spec = _make_slide(LayoutType.TITLE_HERO, page=1, title="Cover")
        renderer.render([slide_spec])
        slide = renderer.prs.slides[0]
        # Count text boxes — none should contain a page number footer
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert not any(t.strip() == "1" for t in texts)

    def test_closing_no_footer(self, renderer):
        slide_spec = _make_slide(LayoutType.CLOSING, page=10, title="End")
        renderer.render([slide_spec])
        slide = renderer.prs.slides[0]
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert not any(t.strip() == "10" for t in texts)

    def test_content_has_page_number(self):
        theme = _make_theme(footer_text="Confidential")
        r = PPTRenderer(theme)
        slide_spec = _make_slide(LayoutType.BULLET_WITH_ICONS, page=7, title="X",
                                 bullets=["A"])
        r.render([slide_spec])
        slide = r.prs.slides[0]
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        assert any("7" in t for t in texts)
        assert any("Confidential" in t for t in texts)


# ===========================================================================
# Test: Image handling (with real temp image)
# ===========================================================================


class TestImageHandling:
    @pytest.fixture
    def tiny_image(self, tmp_path):
        """Create a 10x10 red PNG for testing."""
        try:
            from PIL import Image
        except ImportError:
            pytest.skip("Pillow not installed")
        img = Image.new("RGB", (10, 10), color="red")
        path = str(tmp_path / "test.png")
        img.save(path)
        return path

    def test_text_left_image_right_with_real_image(self, renderer, tiny_image):
        slide = _make_slide(
            LayoutType.TEXT_LEFT_IMAGE_RIGHT,
            title="With Image",
            bullets=["Point"],
            image_path=tiny_image,
        )
        prs = _render_and_reload(renderer, [slide])
        assert len(prs.slides) == 1

    def test_title_hero_with_background_image(self, renderer, tiny_image):
        slide = _make_slide(
            LayoutType.TITLE_HERO,
            title="Hero",
            subtitle="Sub",
            image_path=tiny_image,
        )
        prs = _render_and_reload(renderer, [slide])
        assert len(prs.slides) == 1

    def test_full_image_overlay_with_image(self, renderer, tiny_image):
        slide = _make_slide(
            LayoutType.FULL_IMAGE_OVERLAY,
            title="Overlay",
            image_path=tiny_image,
        )
        prs = _render_and_reload(renderer, [slide])
        assert len(prs.slides) == 1

    def test_missing_image_uses_placeholder(self, renderer):
        """Non-existent image path should not crash — placeholder used."""
        slide = _make_slide(
            LayoutType.TEXT_LEFT_IMAGE_RIGHT,
            title="No Image",
            bullets=["Still works"],
            image_path="/nonexistent/image.png",
        )
        prs = _render_and_reload(renderer, [slide])
        assert len(prs.slides) == 1


# ===========================================================================
# Test: Save
# ===========================================================================


class TestSave:
    def test_save_creates_file(self, renderer):
        renderer.render([_make_slide(LayoutType.CLOSING, title="End")])
        with tempfile.TemporaryDirectory() as tmpdir:
            path = os.path.join(tmpdir, "out.pptx")
            result = renderer.save(path)
            assert os.path.isfile(result)
            assert result == os.path.abspath(path)

    def test_save_creates_parent_dirs(self, renderer):
        renderer.render([_make_slide(LayoutType.CLOSING, title="End")])
        with tempfile.TemporaryDirectory() as tmpdir:
            path = os.path.join(tmpdir, "sub", "dir", "out.pptx")
            result = renderer.save(path)
            assert os.path.isfile(result)


# ===========================================================================
# Test: Different themes
# ===========================================================================


class TestThemeVariations:
    def test_dark_theme(self):
        theme = _make_theme(
            colors=_make_color_scheme(
                primary="#E0E0E0",
                secondary="#BB86FC",
                accent="#03DAC6",
                text="#E0E0E0",
                background="#121212",
            ),
        )
        r = PPTRenderer(theme)
        slides = [
            _make_slide(LayoutType.TITLE_HERO, title="Dark Mode"),
            _make_slide(LayoutType.BULLET_WITH_ICONS, title="Points",
                        bullets=["A: desc", "B: desc"]),
        ]
        prs = _render_and_reload(r, slides)
        assert len(prs.slides) == 2

    def test_creative_theme(self):
        theme = _make_theme(
            colors=_make_color_scheme(
                primary="#E65100",
                secondary="#FF6F61",
                accent="#FDD835",
                text="#3E2723",
                background="#FFF8E1",
            ),
        )
        r = PPTRenderer(theme)
        slides = [_make_slide(LayoutType.QUOTE_PAGE, title="Q",
                              body_text="Be bold.", subtitle="Author")]
        prs = _render_and_reload(r, slides)
        assert len(prs.slides) == 1


# ===========================================================================
# Test: Edge cases
# ===========================================================================


class TestEdgeCases:
    def test_empty_slides_list(self, renderer):
        prs = renderer.render([])
        assert len(prs.slides) == 0

    def test_very_long_title(self, renderer):
        long_title = "A" * 50  # max_length=50 in SlideContent
        slide = _make_slide(LayoutType.SECTION_DIVIDER, title=long_title)
        prs = _render_and_reload(renderer, [slide])
        assert len(prs.slides) == 1

    def test_empty_bullets(self, renderer):
        slide = _make_slide(LayoutType.BULLET_WITH_ICONS, title="Empty",
                            bullets=[])
        prs = _render_and_reload(renderer, [slide])
        assert len(prs.slides) == 1

    def test_comparison_odd_bullets(self, renderer):
        """Odd number of bullets should split without error."""
        slide = _make_slide(LayoutType.COMPARISON, title="Odd",
                            subtitle="A vs B", bullets=["X", "Y", "Z"])
        prs = _render_and_reload(renderer, [slide])
        assert len(prs.slides) == 1

    def test_timeline_single_item(self, renderer):
        slide = _make_slide(LayoutType.TIMELINE, title="One",
                            bullets=["2024: Launch"])
        prs = _render_and_reload(renderer, [slide])
        assert len(prs.slides) == 1

    def test_data_highlight_no_value(self, renderer):
        slide = _make_slide(LayoutType.DATA_HIGHLIGHT, title="No Data")
        prs = _render_and_reload(renderer, [slide])
        assert len(prs.slides) == 1

    def test_three_columns_fewer_than_3_items(self, renderer):
        slide = _make_slide(LayoutType.THREE_COLUMNS, title="Two",
                            bullets=["A: desc", "B: desc"])
        prs = _render_and_reload(renderer, [slide])
        assert len(prs.slides) == 1

    def test_comparison_no_subtitle(self, renderer):
        """Comparison without 'A vs B' subtitle should use defaults."""
        slide = _make_slide(LayoutType.COMPARISON, title="Compare",
                            bullets=["X", "Y", "Z", "W"])
        prs = _render_and_reload(renderer, [slide])
        assert len(prs.slides) == 1
