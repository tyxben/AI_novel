"""Tests for PPT template generation and template-based rendering."""

from __future__ import annotations

import os
import tempfile

import pytest
from pptx import Presentation
from pptx.util import Inches

from src.ppt.create_templates import (
    LAYOUT_BLANK,
    LAYOUT_CLOSING,
    LAYOUT_COMPARISON,
    LAYOUT_DATA_HIGHLIGHT,
    LAYOUT_FULL_IMAGE,
    LAYOUT_IMG_LEFT_TEXT_RIGHT,
    LAYOUT_QUOTE,
    LAYOUT_SECTION,
    LAYOUT_TEXT_LEFT_IMG_RIGHT,
    LAYOUT_THREE_COLUMNS,
    LAYOUT_TITLE,
    SLIDE_HEIGHT,
    SLIDE_WIDTH,
    TEMPLATES_DIR,
    THEMES,
    TemplateTheme,
    create_all_templates,
    create_template,
)
from src.ppt.models import (
    ColorScheme,
    DecorationSpec,
    FontSpec,
    LayoutType,
    SlideContent,
    SlideDesign,
    SlideSpec,
    ThemeConfig,
)
from src.ppt.ppt_renderer import PPTRenderer


# ---------------------------------------------------------------------------
# Template generation tests
# ---------------------------------------------------------------------------


class TestTemplateTheme:
    def test_all_themes_defined(self):
        assert len(THEMES) == 5
        names = {t.name for t in THEMES}
        assert names == {"modern", "business", "creative", "tech", "education"}

    def test_theme_has_required_fields(self):
        for t in THEMES:
            assert t.name
            assert t.display_name
            assert t.primary.startswith("#")
            assert t.accent.startswith("#")
            assert t.bg.startswith("#")

    def test_tech_is_dark_mode(self):
        tech = next(t for t in THEMES if t.name == "tech")
        assert tech.dark_mode is True
        assert tech.bg == "#0D1117"

    def test_secondary_auto_generated(self):
        t = TemplateTheme(
            name="test",
            display_name="Test",
            primary="#FF0000",
            accent="#00FF00",
        )
        # secondary should be auto-generated (muted version of primary)
        assert t.secondary.startswith("#")
        assert t.secondary != ""


class TestCreateTemplate:
    def test_creates_single_template(self, tmp_path, monkeypatch):
        import src.ppt.create_templates as mod
        monkeypatch.setattr(mod, "TEMPLATES_DIR", tmp_path)

        theme = THEMES[0]  # modern
        path = create_template(theme)
        assert path.exists()
        assert path.suffix == ".pptx"
        assert path.name == "modern.pptx"

    def test_template_has_11_slides(self, tmp_path, monkeypatch):
        import src.ppt.create_templates as mod
        monkeypatch.setattr(mod, "TEMPLATES_DIR", tmp_path)

        theme = THEMES[0]
        path = create_template(theme)
        prs = Presentation(str(path))
        assert len(prs.slides) == 11

    def test_template_slide_dimensions(self, tmp_path, monkeypatch):
        import src.ppt.create_templates as mod
        monkeypatch.setattr(mod, "TEMPLATES_DIR", tmp_path)

        theme = THEMES[0]
        path = create_template(theme)
        prs = Presentation(str(path))
        assert prs.slide_width == SLIDE_WIDTH
        assert prs.slide_height == SLIDE_HEIGHT

    @pytest.mark.parametrize("theme", THEMES, ids=[t.name for t in THEMES])
    def test_all_themes_generate(self, theme, tmp_path, monkeypatch):
        import src.ppt.create_templates as mod
        monkeypatch.setattr(mod, "TEMPLATES_DIR", tmp_path)

        path = create_template(theme)
        assert path.exists()
        prs = Presentation(str(path))
        assert len(prs.slides) == 11


class TestCreateAllTemplates:
    def test_creates_all_five(self, tmp_path, monkeypatch):
        import src.ppt.create_templates as mod
        monkeypatch.setattr(mod, "TEMPLATES_DIR", tmp_path)

        paths = create_all_templates()
        assert len(paths) == 5
        for p in paths:
            assert p.exists()

    def test_output_filenames(self, tmp_path, monkeypatch):
        import src.ppt.create_templates as mod
        monkeypatch.setattr(mod, "TEMPLATES_DIR", tmp_path)

        paths = create_all_templates()
        names = {p.name for p in paths}
        assert names == {
            "modern.pptx",
            "business.pptx",
            "creative.pptx",
            "tech.pptx",
            "education.pptx",
        }


class TestTemplateSlideContent:
    """Verify that template slides contain the expected placeholder text."""

    @pytest.fixture
    def modern_prs(self, tmp_path, monkeypatch):
        import src.ppt.create_templates as mod
        monkeypatch.setattr(mod, "TEMPLATES_DIR", tmp_path)
        path = create_template(THEMES[0])
        return Presentation(str(path))

    def _get_texts(self, slide):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            texts.append(run.text.strip())
        return texts

    def test_title_slide_has_placeholders(self, modern_prs):
        texts = self._get_texts(modern_prs.slides[LAYOUT_TITLE])
        assert any("{title}" in t for t in texts)
        assert any("{subtitle}" in t for t in texts)

    def test_section_slide_has_title_placeholder(self, modern_prs):
        texts = self._get_texts(modern_prs.slides[LAYOUT_SECTION])
        assert any("{title}" in t for t in texts)

    def test_quote_slide_has_placeholders(self, modern_prs):
        texts = self._get_texts(modern_prs.slides[LAYOUT_QUOTE])
        assert any("{quote}" in t for t in texts)
        assert any("{author}" in t for t in texts)

    def test_data_highlight_has_placeholders(self, modern_prs):
        texts = self._get_texts(modern_prs.slides[LAYOUT_DATA_HIGHLIGHT])
        assert any("{data_value}" in t for t in texts)
        assert any("{data_label}" in t for t in texts)

    def test_closing_slide_has_placeholders(self, modern_prs):
        texts = self._get_texts(modern_prs.slides[LAYOUT_CLOSING])
        assert any("{title}" in t for t in texts)
        assert any("{contact}" in t for t in texts)

    def test_comparison_has_left_right(self, modern_prs):
        texts = self._get_texts(modern_prs.slides[LAYOUT_COMPARISON])
        assert any("{left_title}" in t for t in texts)
        assert any("{right_title}" in t for t in texts)

    def test_blank_slide_is_empty(self, modern_prs):
        texts = self._get_texts(modern_prs.slides[LAYOUT_BLANK])
        assert len(texts) == 0


# ---------------------------------------------------------------------------
# Template-based rendering integration tests
# ---------------------------------------------------------------------------


def _make_color_scheme(**overrides) -> ColorScheme:
    defaults = dict(
        primary="#2D3436",
        secondary="#636E72",
        accent="#0984E3",
        text="#2D3436",
        background="#FFFFFF",
    )
    defaults.update(overrides)
    return ColorScheme(**defaults)


def _make_font(size=20, bold=False, color="#424242") -> FontSpec:
    return FontSpec(size=size, bold=bold, color=color, family="Arial")


def _make_design(layout: LayoutType) -> SlideDesign:
    return SlideDesign(
        layout=layout,
        colors=_make_color_scheme(),
        title_font=_make_font(32, True, "#2D3436"),
        body_font=_make_font(18, False, "#424242"),
        note_font=_make_font(12, False, "#757575"),
        decoration=DecorationSpec(),
    )


def _make_modern_theme() -> ThemeConfig:
    return ThemeConfig(
        name="modern",
        display_name="Modern",
        colors=_make_color_scheme(),
        title_font=_make_font(32, True, "#2D3436"),
        body_font=_make_font(18, False, "#424242"),
        note_font=_make_font(12, False, "#757575"),
    )


class TestTemplateRendering:
    """Test that the renderer uses templates when available."""

    def test_renderer_loads_template_for_known_theme(self):
        """modern.pptx exists in templates dir, so renderer should load it."""
        theme = _make_modern_theme()
        r = PPTRenderer(theme)
        assert r._template is not None
        assert len(r._template.slides) == 11

    def test_renderer_no_template_for_unknown_theme(self):
        theme = _make_modern_theme()
        theme.name = "nonexistent_theme_xyz"
        r = PPTRenderer(theme)
        assert r._template is None

    def test_template_title_text_replaced(self):
        """When rendering with template, placeholder text should be replaced."""
        theme = _make_modern_theme()
        r = PPTRenderer(theme)

        slide_spec = SlideSpec(
            page_number=1,
            content=SlideContent(
                title="My Presentation",
                subtitle="By Author",
            ),
            design=_make_design(LayoutType.TITLE_HERO),
        )
        r.render([slide_spec])

        slide = r.prs.slides[0]
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        texts.append(run.text)

        assert any("My Presentation" in t for t in texts)
        assert any("By Author" in t for t in texts)
        # Placeholder should be gone
        assert not any("{title}" in t for t in texts)

    def test_template_section_text_replaced(self):
        theme = _make_modern_theme()
        r = PPTRenderer(theme)

        slide_spec = SlideSpec(
            page_number=2,
            content=SlideContent(title="Chapter One"),
            design=_make_design(LayoutType.SECTION_DIVIDER),
        )
        r.render([slide_spec])

        slide = r.prs.slides[0]
        texts = [
            run.text
            for shape in slide.shapes if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            for run in para.runs
        ]
        assert any("Chapter One" in t for t in texts)

    def test_fallback_to_code_for_timeline(self):
        """Timeline always uses code-based rendering (not template)."""
        theme = _make_modern_theme()
        r = PPTRenderer(theme)

        slide_spec = SlideSpec(
            page_number=1,
            content=SlideContent(
                title="Timeline",
                bullet_points=["2020: Start", "2021: Grow"],
            ),
            design=_make_design(LayoutType.TIMELINE),
        )
        prs = r.render([slide_spec])
        assert len(prs.slides) == 1

    def test_fallback_to_code_for_bullet_icons(self):
        """Bullet with icons always uses code-based rendering."""
        theme = _make_modern_theme()
        r = PPTRenderer(theme)

        slide_spec = SlideSpec(
            page_number=1,
            content=SlideContent(
                title="Features",
                bullet_points=["Speed: Fast", "Quality: High"],
            ),
            design=_make_design(LayoutType.BULLET_WITH_ICONS),
        )
        prs = r.render([slide_spec])
        assert len(prs.slides) == 1

    def test_fallback_rendering_without_template(self):
        """With unknown theme name, all layouts fall back to code."""
        theme = _make_modern_theme()
        theme.name = "nonexistent"
        r = PPTRenderer(theme)
        assert r._template is None

        slides = [
            SlideSpec(
                page_number=1,
                content=SlideContent(title="Cover", subtitle="Sub"),
                design=_make_design(LayoutType.TITLE_HERO),
            ),
            SlideSpec(
                page_number=2,
                content=SlideContent(title="Quote", body_text="Be bold."),
                design=_make_design(LayoutType.QUOTE_PAGE),
            ),
        ]
        prs = r.render(slides)
        assert len(prs.slides) == 2

    def test_template_data_highlight_replaced(self):
        theme = _make_modern_theme()
        r = PPTRenderer(theme)

        slide_spec = SlideSpec(
            page_number=1,
            content=SlideContent(
                title="Growth",
                data_value="42%",
                data_label="Year over Year",
            ),
            design=_make_design(LayoutType.DATA_HIGHLIGHT),
        )
        r.render([slide_spec])

        slide = r.prs.slides[0]
        texts = [
            run.text
            for shape in slide.shapes if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            for run in para.runs
        ]
        assert any("42%" in t for t in texts)
        assert any("Year over Year" in t for t in texts)

    def test_template_comparison_replaced(self):
        theme = _make_modern_theme()
        r = PPTRenderer(theme)

        slide_spec = SlideSpec(
            page_number=1,
            content=SlideContent(
                title="Compare",
                subtitle="Plan A vs Plan B",
                bullet_points=["Cheap", "Simple", "Powerful", "Scalable"],
            ),
            design=_make_design(LayoutType.COMPARISON),
        )
        r.render([slide_spec])

        slide = r.prs.slides[0]
        texts = [
            run.text
            for shape in slide.shapes if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            for run in para.runs
        ]
        assert any("Plan A" in t for t in texts)
        assert any("Plan B" in t for t in texts)

    def test_template_closing_replaced(self):
        theme = _make_modern_theme()
        r = PPTRenderer(theme)

        slide_spec = SlideSpec(
            page_number=1,
            content=SlideContent(
                title="Thank You!",
                subtitle="Questions?",
                body_text="me@example.com",
            ),
            design=_make_design(LayoutType.CLOSING),
        )
        r.render([slide_spec])

        slide = r.prs.slides[0]
        texts = [
            run.text
            for shape in slide.shapes if shape.has_text_frame
            for para in shape.text_frame.paragraphs
            for run in para.runs
        ]
        assert any("Thank You!" in t for t in texts)
        assert any("me@example.com" in t for t in texts)

    def test_save_and_reload_with_template(self):
        """Full round-trip: render with template, save, reload."""
        theme = _make_modern_theme()
        r = PPTRenderer(theme)

        slides = [
            SlideSpec(
                page_number=i + 1,
                content=SlideContent(title=f"Slide {i + 1}"),
                design=_make_design(layout),
            )
            for i, layout in enumerate([
                LayoutType.TITLE_HERO,
                LayoutType.SECTION_DIVIDER,
                LayoutType.THREE_COLUMNS,
                LayoutType.CLOSING,
            ])
        ]
        r.render(slides)

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            path = f.name
        try:
            r.save(path)
            reloaded = Presentation(path)
            assert len(reloaded.slides) == 4
        finally:
            os.unlink(path)


class TestBuildReplacements:
    def test_basic_replacements(self):
        spec = SlideSpec(
            page_number=1,
            content=SlideContent(
                title="Title",
                subtitle="Sub",
                body_text="Body",
            ),
            design=_make_design(LayoutType.TITLE_HERO),
        )
        r = PPTRenderer._build_replacements(spec)
        assert r["title"] == "Title"
        assert r["subtitle"] == "Sub"
        assert r["meta"] == "Body"

    def test_bullet_points_in_content(self):
        spec = SlideSpec(
            page_number=1,
            content=SlideContent(
                title="T",
                bullet_points=["A", "B", "C"],
            ),
            design=_make_design(LayoutType.TEXT_LEFT_IMAGE_RIGHT),
        )
        r = PPTRenderer._build_replacements(spec)
        assert "\u2022  A" in r["content"]
        assert "\u2022  B" in r["content"]

    def test_comparison_vs_parsing(self):
        spec = SlideSpec(
            page_number=1,
            content=SlideContent(
                title="T",
                subtitle="A vs B",
                bullet_points=["X", "Y", "Z", "W"],
            ),
            design=_make_design(LayoutType.COMPARISON),
        )
        r = PPTRenderer._build_replacements(spec)
        assert r["left_title"] == "A"
        assert r["right_title"] == "B"
        assert "\u2022  X" in r["left_content"]
        assert "\u2022  Z" in r["right_content"]

    def test_three_columns_splitting(self):
        spec = SlideSpec(
            page_number=1,
            content=SlideContent(
                title="T",
                bullet_points=["Speed: Fast", "Quality: High"],
            ),
            design=_make_design(LayoutType.THREE_COLUMNS),
        )
        r = PPTRenderer._build_replacements(spec)
        assert r["col1_title"] == "Speed"
        assert r["col1_desc"] == "Fast"
        assert r["col2_title"] == "Quality"
        assert r["col2_desc"] == "High"
        assert r["col3_title"] == ""

    def test_empty_content_no_crash(self):
        spec = SlideSpec(
            page_number=1,
            content=SlideContent(title=""),
            design=_make_design(LayoutType.SECTION_DIVIDER),
        )
        r = PPTRenderer._build_replacements(spec)
        assert r["title"] == ""
        assert r["subtitle"] == ""

    def test_default_comparison_labels(self):
        """Without 'vs' in subtitle, defaults to 方案 A / 方案 B."""
        spec = SlideSpec(
            page_number=1,
            content=SlideContent(title="T"),
            design=_make_design(LayoutType.COMPARISON),
        )
        r = PPTRenderer._build_replacements(spec)
        assert r["left_title"] == "\u65B9\u6848 A"
        assert r["right_title"] == "\u65B9\u6848 B"
