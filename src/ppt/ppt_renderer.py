"""PPT 渲染器 - 阶段6：使用 python-pptx 渲染最终 .pptx 文件

核心设计原则：精美、有人味、不模板化。
- 大量留白 = 高级感
- 装饰元素微妙点缀，不堆砌
- 字体层级清晰（4级）
- 配色饱和度 60-80%，不用 100% 纯色
"""

from __future__ import annotations

import copy
import logging
import os
from pathlib import Path
from typing import TYPE_CHECKING

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from .models import LayoutType, SlideSpec, ThemeConfig

if TYPE_CHECKING:
    from pptx.shapes.base import BaseShape
    from pptx.slide import Slide

log = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# 16:9 slide dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Margin constants (converted from px-ish values to Inches for 16:9)
MARGIN_LEFT = Inches(0.9)
MARGIN_RIGHT = Inches(0.9)
MARGIN_TOP = Inches(0.7)
MARGIN_BOTTOM = Inches(0.7)

# Usable area
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
CONTENT_HEIGHT = SLIDE_HEIGHT - MARGIN_TOP - MARGIN_BOTTOM


class PPTRenderer:
    """PPT 渲染引擎，将结构化 SlideSpec 列表渲染为 .pptx 文件。

    支持模板模式：如果 ``templates/<theme_name>.pptx`` 存在，则从模板中
    复制预制的样板页（包含装饰元素），并替换占位文本/插入图片。
    当模板不存在时，退回到纯代码绘制模式。

    Usage::

        renderer = PPTRenderer(theme_config)
        prs = renderer.render(slides)
        renderer.save("output.pptx")
    """

    SLIDE_WIDTH = SLIDE_WIDTH
    SLIDE_HEIGHT = SLIDE_HEIGHT
    TEMPLATES_DIR: Path = Path(__file__).parent / "templates"

    # Mapping from LayoutType to template slide index
    _TEMPLATE_SLIDE_INDEX: dict[LayoutType, int] = {
        LayoutType.TITLE_HERO: 0,
        LayoutType.SECTION_DIVIDER: 1,
        LayoutType.TEXT_LEFT_IMAGE_RIGHT: 2,
        LayoutType.IMAGE_LEFT_TEXT_RIGHT: 3,
        LayoutType.FULL_IMAGE_OVERLAY: 4,
        LayoutType.THREE_COLUMNS: 5,
        LayoutType.QUOTE_PAGE: 6,
        LayoutType.DATA_HIGHLIGHT: 7,
        LayoutType.TIMELINE: 8,         # uses blank template slide
        LayoutType.BULLET_WITH_ICONS: 8,  # uses blank template slide
        LayoutType.COMPARISON: 10,
        LayoutType.CLOSING: 9,
    }

    def __init__(self, theme: ThemeConfig) -> None:
        self.theme = theme
        self._template: Presentation | None = None

        template_path = self.TEMPLATES_DIR / f"{theme.name}.pptx"
        if template_path.exists():
            try:
                self._template = Presentation(str(template_path))
                log.debug("Loaded template: %s", template_path)
            except Exception:
                log.warning("Failed to load template %s, falling back to code rendering",
                            template_path, exc_info=True)
                self._template = None

        self.prs = Presentation()
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def render(self, slides: list[SlideSpec]) -> Presentation:
        """Render all slides and return the Presentation object."""
        for spec in slides:
            self._render_slide(spec)
        return self.prs

    def save(self, path: str) -> str:
        """Save the presentation to *path* and return the absolute path."""
        os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
        self.prs.save(path)
        return os.path.abspath(path)

    # ------------------------------------------------------------------
    # Template cloning helpers
    # ------------------------------------------------------------------

    def _clone_template_slide(self, layout: LayoutType) -> Slide | None:
        """Clone a slide from the template file for the given layout.

        Returns ``None`` if the template is unavailable or the layout index
        is out of range, signaling the caller to fall back to code-based
        rendering.
        """
        if self._template is None:
            return None

        idx = self._TEMPLATE_SLIDE_INDEX.get(layout)
        if idx is None or idx >= len(self._template.slides):
            return None

        try:
            src_slide = self._template.slides[idx]
            # Add a blank slide in the output presentation
            dest_slide = self._add_blank_slide()
            # Copy background
            bg_elem = src_slide.background._element
            dest_slide.background._element.getparent().replace(
                dest_slide.background._element,
                copy.deepcopy(bg_elem),
            )
            # Copy all shapes via XML deep-copy
            for shape in src_slide.shapes:
                el = copy.deepcopy(shape._element)
                dest_slide.shapes._spTree.append(el)
            return dest_slide
        except Exception:
            log.debug("Failed to clone template slide for %s, falling back", layout,
                       exc_info=True)
            return None

    def _replace_placeholder_text(
        self,
        slide: Slide,
        replacements: dict[str, str],
    ) -> None:
        """Replace ``{placeholder}`` tokens in all text frames on the slide.

        Each key in *replacements* should be the placeholder name without
        braces, e.g. ``{"title": "My Title", "subtitle": "Sub"}``.
        """
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for key, val in replacements.items():
                        token = "{" + key + "}"
                        if token in run.text:
                            run.text = run.text.replace(token, val)

    @staticmethod
    def _build_replacements(spec: SlideSpec) -> dict[str, str]:
        """Build a placeholder replacement dict from a SlideSpec.

        Template slides use ``{title}``, ``{subtitle}``, ``{content}`` etc.
        as placeholder tokens inside text frames.  All placeholders are
        populated -- unused ones become empty strings.
        """
        c = spec.content
        bullet_text = "\n".join(f"\u2022  {bp}" for bp in c.bullet_points)
        bullets = c.bullet_points or []

        r: dict[str, str] = {
            "title": c.title or "",
            "subtitle": c.subtitle or "",
            "content": bullet_text or c.body_text or "",
            "meta": c.body_text or "",
            "quote": c.body_text or c.title or "",
            "author": f"\u2014\u2014 {c.subtitle}" if c.subtitle else "",
            "data_value": c.data_value or "",
            "data_label": c.data_label or "",
            "description": c.body_text or (
                c.bullet_points[0] if c.bullet_points else ""
            ),
            "contact": c.body_text or (
                " | ".join(c.bullet_points) if c.bullet_points else ""
            ),
        }

        # Three columns
        for i in range(3):
            if i < len(bullets):
                text = bullets[i]
                parts = text.split(":", 1) if ":" in text else [text, ""]
                r[f"col{i+1}_title"] = parts[0].strip()
                r[f"col{i+1}_desc"] = parts[1].strip() if len(parts) > 1 else ""
            else:
                r[f"col{i+1}_title"] = ""
                r[f"col{i+1}_desc"] = ""

        # Comparison
        mid = len(bullets) // 2
        sub = c.subtitle or ""
        if " vs " in sub.lower():
            parts = sub.split(" vs ", 1)
            r["left_title"] = parts[0].strip()
            r["right_title"] = parts[1].strip() if len(parts) > 1 else ""
        elif "vs" in sub.lower():
            parts = sub.lower().split("vs", 1)
            r["left_title"] = parts[0].strip()
            r["right_title"] = parts[1].strip() if len(parts) > 1 else ""
        else:
            r["left_title"] = "\u65B9\u6848 A"
            r["right_title"] = "\u65B9\u6848 B"

        r["left_content"] = "\n".join(f"\u2022  {b}" for b in bullets[:mid])
        r["right_content"] = "\n".join(f"\u2022  {b}" for b in bullets[mid:])
        return r

    # ------------------------------------------------------------------
    # Slide dispatch
    # ------------------------------------------------------------------

    _LAYOUT_RENDERERS: dict[LayoutType, str] = {
        LayoutType.TITLE_HERO: "_render_title_hero",
        LayoutType.SECTION_DIVIDER: "_render_section_divider",
        LayoutType.TEXT_LEFT_IMAGE_RIGHT: "_render_text_left_image_right",
        LayoutType.IMAGE_LEFT_TEXT_RIGHT: "_render_image_left_text_right",
        LayoutType.FULL_IMAGE_OVERLAY: "_render_full_image_overlay",
        LayoutType.THREE_COLUMNS: "_render_three_columns",
        LayoutType.QUOTE_PAGE: "_render_quote_page",
        LayoutType.DATA_HIGHLIGHT: "_render_data_highlight",
        LayoutType.TIMELINE: "_render_timeline",
        LayoutType.BULLET_WITH_ICONS: "_render_bullet_with_icons",
        LayoutType.COMPARISON: "_render_comparison",
        LayoutType.CLOSING: "_render_closing",
    }

    # Layouts where template clone + text replacement is sufficient.
    # Timeline and bullet_with_icons need dynamic shape generation,
    # so they always use the code-based renderer.
    _TEMPLATE_ONLY_LAYOUTS: set[LayoutType] = {
        LayoutType.TITLE_HERO,
        LayoutType.SECTION_DIVIDER,
        LayoutType.THREE_COLUMNS,
        LayoutType.QUOTE_PAGE,
        LayoutType.DATA_HIGHLIGHT,
        LayoutType.CLOSING,
        LayoutType.COMPARISON,
        LayoutType.FULL_IMAGE_OVERLAY,
        LayoutType.TEXT_LEFT_IMAGE_RIGHT,
        LayoutType.IMAGE_LEFT_TEXT_RIGHT,
    }

    def _render_slide(self, spec: SlideSpec) -> None:
        method_name = self._LAYOUT_RENDERERS.get(spec.design.layout)
        if method_name is None:
            return

        layout = spec.design.layout
        slide: Slide | None = None

        # Try template-based rendering for supported layouts
        if layout in self._TEMPLATE_ONLY_LAYOUTS:
            slide = self._clone_template_slide(layout)
            if slide is not None:
                replacements = self._build_replacements(spec)
                self._replace_placeholder_text(slide, replacements)
                # Handle image replacement on template slides
                self._handle_template_images(slide, spec)

        if slide is None:
            # Fallback: draw from scratch on a blank slide
            method = getattr(self, method_name)
            slide = self._add_blank_slide()
            bg_color = spec.design.colors.background
            self._set_slide_background(slide, bg_color)
            method(slide, spec)
        # Speaker notes
        if spec.content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = spec.content.speaker_notes
        # Footer / page number (skip for title and closing)
        if spec.design.layout not in (LayoutType.TITLE_HERO, LayoutType.CLOSING):
            self._add_footer(
                slide,
                text=self.theme.footer_text or "",
                page_number=spec.page_number,
            )

    def _handle_template_images(self, slide: Slide, spec: SlideSpec) -> None:
        """Insert images into a template-cloned slide when applicable."""
        if not spec.image_path or not os.path.isfile(spec.image_path):
            return

        layout = spec.design.layout

        if layout in (LayoutType.TEXT_LEFT_IMAGE_RIGHT, LayoutType.IMAGE_LEFT_TEXT_RIGHT):
            # Replace the rounded-rect placeholder with the actual image
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            for shape in list(slide.shapes):
                if (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                        and shape.width > Inches(3)
                        and shape.height > Inches(3)):
                    try:
                        if shape.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
                            left, top, w, h = shape.left, shape.top, shape.width, shape.height
                            shape._element.getparent().remove(shape._element)
                            self._add_image(slide, spec.image_path, left, top, w, h)
                            break
                    except Exception:
                        pass

        elif layout == LayoutType.FULL_IMAGE_OVERLAY:
            pic = self._add_image(
                slide, spec.image_path, Emu(0), Emu(0),
                self.SLIDE_WIDTH, self.SLIDE_HEIGHT,
            )
            # Move image to back
            sp_tree = slide.shapes._spTree
            pic_elem = pic._element
            sp_tree.remove(pic_elem)
            sp_tree.insert(2, pic_elem)

        elif layout == LayoutType.TITLE_HERO:
            pic = self._add_image(
                slide, spec.image_path, Emu(0), Emu(0),
                self.SLIDE_WIDTH, self.SLIDE_HEIGHT,
            )
            sp_tree = slide.shapes._spTree
            pic_elem = pic._element
            sp_tree.remove(pic_elem)
            sp_tree.insert(2, pic_elem)
            self._add_overlay(slide, color="#000000", opacity=0.50)

    def _add_blank_slide(self) -> Slide:
        """Add a blank slide (layout index 6 = Blank)."""
        layout = self.prs.slide_layouts[6]
        return self.prs.slides.add_slide(layout)

    # ==================================================================
    # Layout Renderers
    # ==================================================================

    def _render_title_hero(self, slide: Slide, spec: SlideSpec) -> None:
        """封面页: 大标题 + 副标题 + 底部渐变条 + 可选背景图."""
        colors = spec.design.colors

        # Background image or gradient accent bar at bottom
        if spec.image_path and os.path.isfile(spec.image_path):
            self._add_image(slide, spec.image_path, Emu(0), Emu(0),
                            self.SLIDE_WIDTH, self.SLIDE_HEIGHT)
            # Dark overlay for readability
            self._add_overlay(slide, color="#000000", opacity=0.50)
            title_color = "#FFFFFF"
            subtitle_color = "#E0E0E0"
            meta_color = "#BDBDBD"
        else:
            # Accent gradient bar at bottom (15% height)
            bar_height = Inches(1.125)  # 15% of 7.5
            self._add_gradient_rect(
                slide,
                left=Emu(0),
                top=self.SLIDE_HEIGHT - bar_height,
                width=self.SLIDE_WIDTH,
                height=bar_height,
                color1=colors.primary,
                color2=colors.accent,
            )
            title_color = colors.primary
            subtitle_color = colors.text
            meta_color = "#9E9E9E"

        # Title — vertically centered, slight upward offset
        title_top = Inches(2.2)
        self._add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=title_top,
            width=CONTENT_WIDTH,
            height=Inches(1.2),
            text=spec.content.title,
            font_size=42,
            bold=True,
            color=title_color,
            alignment=PP_ALIGN.CENTER,
            font_name=spec.design.title_font.family,
        )

        # Subtitle
        if spec.content.subtitle:
            self._add_text_box(
                slide,
                left=MARGIN_LEFT,
                top=title_top + Inches(1.3),
                width=CONTENT_WIDTH,
                height=Inches(0.8),
                text=spec.content.subtitle,
                font_size=22,
                bold=False,
                color=subtitle_color,
                alignment=PP_ALIGN.CENTER,
                font_name=spec.design.body_font.family,
            )

        # Thin divider line
        line_top = title_top + Inches(2.3)
        line_width = Inches(2.0)
        line_left = (self.SLIDE_WIDTH - line_width) // 2
        self._add_accent_line(slide, line_left, line_top, line_width,
                              color=colors.accent)

        # Date / author metadata
        if spec.content.body_text:
            self._add_text_box(
                slide,
                left=MARGIN_LEFT,
                top=line_top + Inches(0.35),
                width=CONTENT_WIDTH,
                height=Inches(0.5),
                text=spec.content.body_text,
                font_size=13,
                color=meta_color,
                alignment=PP_ALIGN.CENTER,
                font_name=spec.design.body_font.family,
            )

    def _render_section_divider(self, slide: Slide, spec: SlideSpec) -> None:
        """章节分隔页: 大量留白 + 居中标题 + 上下 accent 细线."""
        colors = spec.design.colors

        # Slightly darker background
        self._set_slide_background(slide, self._lighten_or_darken(
            colors.background, darken=8))

        center_y = Inches(3.0)
        line_width = Inches(2.0)
        line_left = (self.SLIDE_WIDTH - line_width) // 2

        # Top accent line
        self._add_accent_line(slide, line_left, center_y, line_width,
                              color=colors.accent)

        # Title
        self._add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=center_y + Inches(0.3),
            width=CONTENT_WIDTH,
            height=Inches(1.0),
            text=spec.content.title,
            font_size=36,
            bold=True,
            color=colors.primary,
            alignment=PP_ALIGN.CENTER,
            font_name=spec.design.title_font.family,
        )

        # Bottom accent line
        self._add_accent_line(slide, line_left, center_y + Inches(1.5),
                              line_width, color=colors.accent)

    def _render_text_left_image_right(self, slide: Slide, spec: SlideSpec) -> None:
        """左文右图: 左 55% 文字 + 右 40% 图片, 5% 间距."""
        self._render_text_image_split(slide, spec, text_on_left=True)

    def _render_image_left_text_right(self, slide: Slide, spec: SlideSpec) -> None:
        """左图右文: 镜像版."""
        self._render_text_image_split(slide, spec, text_on_left=False)

    def _render_text_image_split(self, slide: Slide, spec: SlideSpec,
                                 *, text_on_left: bool) -> None:
        """Shared implementation for text+image split layouts."""
        colors = spec.design.colors
        text_width = Inches(6.3)   # ~55%
        image_width = Inches(4.6)  # ~40%
        gap = Inches(0.5)

        if text_on_left:
            text_left = MARGIN_LEFT
            img_left = MARGIN_LEFT + text_width + gap
        else:
            img_left = MARGIN_LEFT
            text_left = MARGIN_LEFT + image_width + gap

        # Title
        self._add_text_box(
            slide,
            left=text_left,
            top=MARGIN_TOP,
            width=text_width,
            height=Inches(0.9),
            text=spec.content.title,
            font_size=28,
            bold=True,
            color=colors.primary,
            font_name=spec.design.title_font.family,
        )

        # Accent line under title
        self._add_accent_line(
            slide, text_left, MARGIN_TOP + Inches(0.95),
            Inches(1.5), color=colors.accent,
        )

        # Bullet points
        bullet_top = MARGIN_TOP + Inches(1.4)
        bullet_text = "\n".join(f"\u2022  {bp}" for bp in spec.content.bullet_points)
        if bullet_text:
            self._add_text_box(
                slide,
                left=text_left,
                top=bullet_top,
                width=text_width,
                height=Inches(4.5),
                text=bullet_text,
                font_size=18,
                color=colors.text,
                line_spacing=1.6,
                font_name=spec.design.body_font.family,
            )

        # Image area
        img_top = MARGIN_TOP + Inches(0.3)
        img_height = Inches(5.5)
        if spec.image_path and os.path.isfile(spec.image_path):
            self._add_image(slide, spec.image_path, img_left, img_top,
                            image_width, img_height)
        else:
            # Placeholder color block with rounded-rect feel
            self._add_shape(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                img_left, img_top, image_width, img_height,
                fill_color=self._muted_color(colors.secondary, 0.15),
            )

    def _render_full_image_overlay(self, slide: Slide, spec: SlideSpec) -> None:
        """全屏图片 + 半透明暗色覆盖 + 白色文字."""
        colors = spec.design.colors

        if spec.image_path and os.path.isfile(spec.image_path):
            self._add_image(slide, spec.image_path, Emu(0), Emu(0),
                            self.SLIDE_WIDTH, self.SLIDE_HEIGHT)
        else:
            # Dark gradient background as fallback
            self._add_gradient_rect(slide, Emu(0), Emu(0),
                                    self.SLIDE_WIDTH, self.SLIDE_HEIGHT,
                                    colors.primary, "#1A1A2E")

        # Dark overlay
        self._add_overlay(slide, color="#000000", opacity=0.55)

        # Title in lower-center area
        title_top = Inches(3.5)
        self._add_text_box(
            slide,
            left=Inches(1.5),
            top=title_top,
            width=Inches(10.333),
            height=Inches(1.2),
            text=spec.content.title,
            font_size=38,
            bold=True,
            color="#FFFFFF",
            alignment=PP_ALIGN.CENTER,
            font_name=spec.design.title_font.family,
        )

        # Short body text
        body = spec.content.body_text or (
            "\n".join(spec.content.bullet_points[:2]) if spec.content.bullet_points else ""
        )
        if body:
            self._add_text_box(
                slide,
                left=Inches(2.0),
                top=title_top + Inches(1.4),
                width=Inches(9.333),
                height=Inches(1.0),
                text=body,
                font_size=20,
                color="#E0E0E0",
                alignment=PP_ALIGN.CENTER,
                font_name=spec.design.body_font.family,
            )

    def _render_three_columns(self, slide: Slide, spec: SlideSpec) -> None:
        """三栏并列: 页面标题 + 3列内容."""
        colors = spec.design.colors

        # Page title
        self._add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=MARGIN_TOP,
            width=CONTENT_WIDTH,
            height=Inches(0.9),
            text=spec.content.title,
            font_size=28,
            bold=True,
            color=colors.primary,
            alignment=PP_ALIGN.CENTER,
            font_name=spec.design.title_font.family,
        )

        # Thin full-width divider
        self._add_accent_line(
            slide,
            left=Inches(2.0),
            top=MARGIN_TOP + Inches(1.0),
            width=Inches(9.333),
            color=colors.secondary,
        )

        # Three columns
        col_width = Inches(3.4)
        col_gap = Inches(0.37)
        col_start_left = MARGIN_LEFT + Inches(0.35)
        col_top = MARGIN_TOP + Inches(1.5)
        col_height = Inches(4.5)

        bullets = spec.content.bullet_points or []
        # Distribute bullets evenly across 3 columns (at most 3 items)
        for i in range(3):
            col_left = col_start_left + i * (col_width + col_gap)

            # Decorative circle at top of column
            circle_size = Inches(0.55)
            circle_left = col_left + (col_width - circle_size) // 2
            self._add_shape(
                slide, MSO_SHAPE.OVAL,
                circle_left, col_top, circle_size, circle_size,
                fill_color=colors.accent if i == 0 else (
                    colors.secondary if i == 1 else colors.primary),
            )

            if i < len(bullets):
                # Column text
                text = bullets[i]
                # Split on first sentence boundary or use full text
                parts = text.split(":", 1) if ":" in text else [text, ""]
                heading = parts[0].strip()
                desc = parts[1].strip() if len(parts) > 1 and parts[1].strip() else ""

                # Column heading
                self._add_text_box(
                    slide,
                    left=col_left,
                    top=col_top + Inches(0.8),
                    width=col_width,
                    height=Inches(0.6),
                    text=heading,
                    font_size=18,
                    bold=True,
                    color=colors.primary,
                    alignment=PP_ALIGN.CENTER,
                    font_name=spec.design.title_font.family,
                )

                if desc:
                    self._add_text_box(
                        slide,
                        left=col_left,
                        top=col_top + Inches(1.5),
                        width=col_width,
                        height=Inches(2.5),
                        text=desc,
                        font_size=15,
                        color=colors.text,
                        alignment=PP_ALIGN.CENTER,
                        line_spacing=1.5,
                        font_name=spec.design.body_font.family,
                    )

    def _render_quote_page(self, slide: Slide, spec: SlideSpec) -> None:
        """引用/金句页: 大引号装饰 + 斜体引文居中 + 作者名."""
        colors = spec.design.colors

        # Large decorative open-quote character
        self._add_text_box(
            slide,
            left=Inches(1.5),
            top=Inches(1.0),
            width=Inches(2.0),
            height=Inches(1.8),
            text="\u201C",
            font_size=120,
            bold=True,
            color=self._muted_color(colors.accent, 0.35),
            alignment=PP_ALIGN.LEFT,
            font_name="Georgia",
        )

        # Quote text
        quote_text = spec.content.body_text or spec.content.title
        self._add_text_box(
            slide,
            left=Inches(2.0),
            top=Inches(2.5),
            width=Inches(9.333),
            height=Inches(2.5),
            text=quote_text,
            font_size=26,
            italic=True,
            color=colors.text,
            alignment=PP_ALIGN.CENTER,
            line_spacing=1.7,
            font_name=spec.design.body_font.family,
        )

        # Author attribution
        author = spec.content.subtitle or ""
        if author:
            self._add_text_box(
                slide,
                left=Inches(2.0),
                top=Inches(5.3),
                width=Inches(9.333),
                height=Inches(0.5),
                text=f"\u2014\u2014 {author}",
                font_size=16,
                color="#9E9E9E",
                alignment=PP_ALIGN.RIGHT,
                font_name=spec.design.body_font.family,
            )

    def _render_data_highlight(self, slide: Slide, spec: SlideSpec) -> None:
        """数据展示页: 超大数字 + 标签 + 描述."""
        colors = spec.design.colors

        # Page title (smaller, top)
        self._add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=MARGIN_TOP,
            width=CONTENT_WIDTH,
            height=Inches(0.8),
            text=spec.content.title,
            font_size=24,
            bold=True,
            color=colors.primary,
            alignment=PP_ALIGN.CENTER,
            font_name=spec.design.title_font.family,
        )

        # Decorative circle behind the number
        circle_size = Inches(3.0)
        circle_left = (self.SLIDE_WIDTH - circle_size) // 2
        circle_top = Inches(2.0)
        self._add_shape(
            slide, MSO_SHAPE.OVAL,
            circle_left, circle_top, circle_size, circle_size,
            fill_color=self._muted_color(colors.accent, 0.12),
        )

        # Big data value
        data_value = spec.content.data_value or ""
        self._add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=Inches(2.3),
            width=CONTENT_WIDTH,
            height=Inches(1.8),
            text=data_value,
            font_size=72,
            bold=True,
            color=colors.accent,
            alignment=PP_ALIGN.CENTER,
            font_name=spec.design.title_font.family,
        )

        # Data label
        data_label = spec.content.data_label or ""
        if data_label:
            self._add_text_box(
                slide,
                left=MARGIN_LEFT,
                top=Inches(4.3),
                width=CONTENT_WIDTH,
                height=Inches(0.7),
                text=data_label,
                font_size=22,
                bold=True,
                color=colors.primary,
                alignment=PP_ALIGN.CENTER,
                font_name=spec.design.body_font.family,
            )

        # Description from bullet points
        desc = spec.content.body_text or (
            spec.content.bullet_points[0] if spec.content.bullet_points else ""
        )
        if desc:
            self._add_text_box(
                slide,
                left=Inches(2.5),
                top=Inches(5.2),
                width=Inches(8.333),
                height=Inches(0.8),
                text=desc,
                font_size=16,
                color="#757575",
                alignment=PP_ALIGN.CENTER,
                font_name=spec.design.body_font.family,
            )

    def _render_timeline(self, slide: Slide, spec: SlideSpec) -> None:
        """时间线: 标题 + 水平线 + 圆形节点 + 标签."""
        colors = spec.design.colors

        # Page title
        self._add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=MARGIN_TOP,
            width=CONTENT_WIDTH,
            height=Inches(0.9),
            text=spec.content.title,
            font_size=28,
            bold=True,
            color=colors.primary,
            alignment=PP_ALIGN.CENTER,
            font_name=spec.design.title_font.family,
        )

        items = spec.content.bullet_points or []
        n = min(len(items), 5)  # max 5 nodes
        if n == 0:
            return

        # Timeline geometry
        line_y = Inches(3.5)
        line_left = Inches(1.5)
        line_right = Inches(11.833)
        line_width = line_right - line_left

        # Horizontal connector line
        self._add_shape(
            slide, MSO_SHAPE.RECTANGLE,
            line_left, line_y, line_width, Inches(0.04),
            fill_color="#BDBDBD",
        )

        # Nodes
        node_size = Inches(0.4)
        spacing = line_width // n if n > 1 else Emu(0)

        for i in range(n):
            if n == 1:
                cx = line_left + line_width // 2
            else:
                cx = line_left + Emu(int(spacing) * i) + spacing // 2

            # Circle node
            self._add_shape(
                slide, MSO_SHAPE.OVAL,
                cx - node_size // 2, line_y - node_size // 2,
                node_size, node_size,
                fill_color=colors.accent,
            )

            # Parse "label: description" or use full text
            text = items[i]
            parts = text.split(":", 1) if ":" in text else [text, ""]
            label = parts[0].strip()
            desc = parts[1].strip() if len(parts) > 1 else ""

            # Label above node
            label_width = Inches(2.0)
            self._add_text_box(
                slide,
                left=cx - label_width // 2,
                top=line_y - Inches(1.1),
                width=label_width,
                height=Inches(0.5),
                text=label,
                font_size=14,
                bold=True,
                color=colors.primary,
                alignment=PP_ALIGN.CENTER,
                font_name=spec.design.title_font.family,
            )

            # Description below node
            if desc:
                self._add_text_box(
                    slide,
                    left=cx - label_width // 2,
                    top=line_y + Inches(0.5),
                    width=label_width,
                    height=Inches(1.5),
                    text=desc,
                    font_size=13,
                    color=colors.text,
                    alignment=PP_ALIGN.CENTER,
                    line_spacing=1.3,
                    font_name=spec.design.body_font.family,
                )

    def _render_bullet_with_icons(self, slide: Slide, spec: SlideSpec) -> None:
        """要点列表 + 彩色圆点图标."""
        colors = spec.design.colors

        # Title
        self._add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=MARGIN_TOP,
            width=CONTENT_WIDTH,
            height=Inches(0.9),
            text=spec.content.title,
            font_size=28,
            bold=True,
            color=colors.primary,
            font_name=spec.design.title_font.family,
        )

        # Accent line
        self._add_accent_line(slide, MARGIN_LEFT, MARGIN_TOP + Inches(0.95),
                              Inches(1.5), color=colors.accent)

        # Bullet items
        items = spec.content.bullet_points or []
        start_y = MARGIN_TOP + Inches(1.5)
        item_height = Inches(1.3)
        circle_size = Inches(0.3)
        text_indent = Inches(0.6)

        # Cycle through accent colors for variety
        dot_colors = [colors.accent, colors.secondary, colors.primary]

        for i, item in enumerate(items[:5]):
            y = start_y + Emu(int(item_height) * i)
            dot_color = dot_colors[i % len(dot_colors)]

            # Colored circle
            self._add_shape(
                slide, MSO_SHAPE.OVAL,
                MARGIN_LEFT, y + Inches(0.08),
                circle_size, circle_size,
                fill_color=dot_color,
            )

            # Parse heading : description
            parts = item.split(":", 1) if ":" in item else [item, ""]
            heading = parts[0].strip()
            desc = parts[1].strip() if len(parts) > 1 and parts[1].strip() else ""

            # Heading (bold)
            self._add_text_box(
                slide,
                left=MARGIN_LEFT + text_indent,
                top=y,
                width=CONTENT_WIDTH - text_indent,
                height=Inches(0.45),
                text=heading,
                font_size=19,
                bold=True,
                color=colors.primary,
                font_name=spec.design.title_font.family,
            )

            # Description
            if desc:
                self._add_text_box(
                    slide,
                    left=MARGIN_LEFT + text_indent,
                    top=y + Inches(0.45),
                    width=CONTENT_WIDTH - text_indent,
                    height=Inches(0.6),
                    text=desc,
                    font_size=16,
                    color=colors.text,
                    line_spacing=1.4,
                    font_name=spec.design.body_font.family,
                )

    def _render_comparison(self, slide: Slide, spec: SlideSpec) -> None:
        """对比页: 左右两栏, 不同 accent 色."""
        colors = spec.design.colors

        # Title
        self._add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=MARGIN_TOP,
            width=CONTENT_WIDTH,
            height=Inches(0.9),
            text=spec.content.title,
            font_size=28,
            bold=True,
            color=colors.primary,
            alignment=PP_ALIGN.CENTER,
            font_name=spec.design.title_font.family,
        )

        # Two columns
        col_width = Inches(5.2)
        col_gap = Inches(1.0)
        col_top = MARGIN_TOP + Inches(1.3)
        col_height = Inches(5.0)
        left_col_left = MARGIN_LEFT + Inches(0.2)
        right_col_left = left_col_left + col_width + col_gap

        bullets = spec.content.bullet_points or []
        mid = len(bullets) // 2

        # Column header backgrounds
        header_h = Inches(0.7)
        for idx, (col_left, col_color) in enumerate([
            (left_col_left, colors.primary),
            (right_col_left, colors.secondary),
        ]):
            # Header bar
            self._add_shape(
                slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                col_left, col_top, col_width, header_h,
                fill_color=col_color,
            )

            # Determine column label from subtitle or generate default
            if spec.content.subtitle and " vs " in spec.content.subtitle.lower():
                parts = spec.content.subtitle.split(" vs ", 1)
                label = parts[0].strip() if idx == 0 else parts[1].strip()
            elif spec.content.subtitle and "vs" in spec.content.subtitle.lower():
                parts = spec.content.subtitle.lower().split("vs", 1)
                label = parts[0].strip() if idx == 0 else parts[1].strip()
            else:
                label = f"\u65B9\u6848 {'A' if idx == 0 else 'B'}"

            self._add_text_box(
                slide,
                left=col_left,
                top=col_top + Inches(0.05),
                width=col_width,
                height=Inches(0.5),
                text=label,
                font_size=20,
                bold=True,
                color="#FFFFFF",
                alignment=PP_ALIGN.CENTER,
                font_name=spec.design.title_font.family,
            )

            # Column bullet points
            col_bullets = bullets[:mid] if idx == 0 else bullets[mid:]
            bullet_text = "\n".join(f"\u2022  {b}" for b in col_bullets)
            if bullet_text:
                self._add_text_box(
                    slide,
                    left=col_left + Inches(0.25),
                    top=col_top + header_h + Inches(0.3),
                    width=col_width - Inches(0.5),
                    height=col_height - header_h - Inches(0.5),
                    text=bullet_text,
                    font_size=17,
                    color=colors.text,
                    line_spacing=1.6,
                    font_name=spec.design.body_font.family,
                )

    def _render_closing(self, slide: Slide, spec: SlideSpec) -> None:
        """结束页: 感谢 + 联系方式 + 底部渐变条."""
        colors = spec.design.colors

        # Bottom gradient bar
        bar_height = Inches(0.8)
        self._add_gradient_rect(
            slide,
            left=Emu(0),
            top=self.SLIDE_HEIGHT - bar_height,
            width=self.SLIDE_WIDTH,
            height=bar_height,
            color1=colors.primary,
            color2=colors.accent,
        )

        # Main text (Thanks!)
        self._add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=Inches(2.5),
            width=CONTENT_WIDTH,
            height=Inches(1.2),
            text=spec.content.title,
            font_size=40,
            bold=True,
            color=colors.primary,
            alignment=PP_ALIGN.CENTER,
            font_name=spec.design.title_font.family,
        )

        # Subtitle
        if spec.content.subtitle:
            self._add_text_box(
                slide,
                left=MARGIN_LEFT,
                top=Inches(3.8),
                width=CONTENT_WIDTH,
                height=Inches(0.7),
                text=spec.content.subtitle,
                font_size=20,
                color=colors.text,
                alignment=PP_ALIGN.CENTER,
                font_name=spec.design.body_font.family,
            )

        # Contact info from body_text or bullet_points
        contact = spec.content.body_text or (
            " | ".join(spec.content.bullet_points) if spec.content.bullet_points else ""
        )
        if contact:
            self._add_text_box(
                slide,
                left=MARGIN_LEFT,
                top=Inches(4.8),
                width=CONTENT_WIDTH,
                height=Inches(0.5),
                text=contact,
                font_size=14,
                color="#9E9E9E",
                alignment=PP_ALIGN.CENTER,
                font_name=spec.design.body_font.family,
            )

    # ==================================================================
    # Utility / decoration methods
    # ==================================================================

    def _add_text_box(
        self,
        slide: Slide,
        left: int | Emu,
        top: int | Emu,
        width: int | Emu,
        height: int | Emu,
        text: str,
        *,
        font_name: str = "Arial",
        font_size: int = 16,
        bold: bool = False,
        italic: bool = False,
        color: str = "#2D3436",
        alignment: PP_ALIGN = PP_ALIGN.LEFT,
        line_spacing: float = 1.5,
        vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    ) -> BaseShape:
        """Add a text box with full typography control."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        # Vertical anchor
        tf.paragraphs[0].alignment = alignment

        # Set text
        p = tf.paragraphs[0]
        p.text = ""
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = self._hex_to_rgb(color)
        p.alignment = alignment

        # Line spacing via XML (pptx doesn't expose line spacing directly)
        pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn("a:lnSpc"), {})
        spcPct = lnSpc.makeelement(
            qn("a:spcPct"),
            {"val": str(int(line_spacing * 100000))},
        )
        lnSpc.append(spcPct)
        pPr.append(lnSpc)

        return txBox

    def _add_image(
        self,
        slide: Slide,
        image_path: str,
        left: int | Emu,
        top: int | Emu,
        width: int | Emu,
        height: int | Emu,
    ) -> BaseShape:
        """Add an image to the slide, scaling to fit the given rectangle."""
        return slide.shapes.add_picture(image_path, left, top, width, height)

    def _add_shape(
        self,
        slide: Slide,
        shape_type: MSO_SHAPE,
        left: int | Emu,
        top: int | Emu,
        width: int | Emu,
        height: int | Emu,
        *,
        fill_color: str | None = None,
        line_color: str | None = None,
        line_width: float | None = None,
    ) -> BaseShape:
        """Add a shape (rectangle, oval, etc.)."""
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)

        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._hex_to_rgb(fill_color)
        else:
            shape.fill.background()

        if line_color:
            shape.line.color.rgb = self._hex_to_rgb(line_color)
            if line_width is not None:
                shape.line.width = Pt(line_width)
        else:
            shape.line.fill.background()

        return shape

    def _add_gradient_rect(
        self,
        slide: Slide,
        left: int | Emu,
        top: int | Emu,
        width: int | Emu,
        height: int | Emu,
        color1: str,
        color2: str,
    ) -> BaseShape:
        """Add a rectangle with a two-stop gradient fill.

        python-pptx doesn't directly support gradient angle on shapes,
        so we use XML manipulation for a left-to-right linear gradient.
        """
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.line.fill.background()

        # Build gradient fill XML
        spPr = shape._element.spPr
        # Remove any existing fill
        for child in list(spPr):
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag in ("solidFill", "gradFill", "noFill", "pattFill"):
                spPr.remove(child)

        gradFill = spPr.makeelement(qn("a:gradFill"), {"flip": "none"})
        gsLst = gradFill.makeelement(qn("a:gsLst"), {})

        for pos, hex_color in [("0", color1), ("100000", color2)]:
            gs = gsLst.makeelement(qn("a:gs"), {"pos": pos})
            srgbClr = gs.makeelement(qn("a:srgbClr"), {"val": hex_color.lstrip("#")})
            gs.append(srgbClr)
            gsLst.append(gs)

        gradFill.append(gsLst)

        # Linear gradient direction (left to right)
        lin = gradFill.makeelement(qn("a:lin"), {"ang": "0", "scaled": "1"})
        gradFill.append(lin)

        spPr.append(gradFill)

        return shape

    def _add_accent_line(
        self,
        slide: Slide,
        left: int | Emu,
        top: int | Emu,
        width: int | Emu,
        color: str | None = None,
    ) -> BaseShape:
        """Add a thin accent-color horizontal line (decorative divider)."""
        line_color = color or self.theme.colors.accent
        return self._add_shape(
            slide, MSO_SHAPE.RECTANGLE,
            left, top, width, Inches(0.03),
            fill_color=line_color,
        )

    def _add_footer(
        self,
        slide: Slide,
        text: str = "",
        page_number: int = 0,
    ) -> None:
        """Add a subtle footer with optional text and page number."""
        footer_y = self.SLIDE_HEIGHT - Inches(0.45)
        footer_parts: list[str] = []
        if text:
            footer_parts.append(text)
        if page_number > 0:
            footer_parts.append(str(page_number))

        footer_str = "  |  ".join(footer_parts) if footer_parts else ""
        if not footer_str:
            return

        self._add_text_box(
            slide,
            left=MARGIN_LEFT,
            top=footer_y,
            width=CONTENT_WIDTH,
            height=Inches(0.35),
            text=footer_str,
            font_size=10,
            color="#BDBDBD",
            alignment=PP_ALIGN.RIGHT,
            font_name=self.theme.body_font.family,
            line_spacing=1.0,
        )

    def _add_overlay(
        self,
        slide: Slide,
        color: str = "#000000",
        opacity: float = 0.6,
    ) -> BaseShape:
        """Add a semi-transparent dark overlay covering the full slide.

        Used on top of full-bleed images to ensure text readability.
        """
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(0), Emu(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT,
        )
        shape.line.fill.background()

        # Solid fill with alpha
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._hex_to_rgb(color)

        # Set transparency via XML (python-pptx doesn't expose fill alpha)
        # Access the underlying lxml element for the shape's spPr > solidFill
        spPr = shape._element.spPr
        solidFill_elem = spPr.find(qn("a:solidFill"))
        if solidFill_elem is not None:
            srgbClr = solidFill_elem.find(qn("a:srgbClr"))
            if srgbClr is not None:
                alpha_val = str(int((1 - opacity) * 100000))
                alpha_elem = srgbClr.makeelement(qn("a:alpha"), {"val": alpha_val})
                srgbClr.append(alpha_elem)

        return shape

    # ==================================================================
    # Color helpers
    # ==================================================================

    @staticmethod
    def _hex_to_rgb(hex_color: str) -> RGBColor:
        """Convert '#RRGGBB' or 'RRGGBB' to RGBColor."""
        h = hex_color.lstrip("#")
        if len(h) != 6:
            h = "2D3436"  # safe fallback
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    @staticmethod
    def _lighten_or_darken(hex_color: str, *, darken: int = 0, lighten: int = 0) -> str:
        """Shift a hex color darker or lighter by the given amount (0-255)."""
        h = hex_color.lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        r = max(0, min(255, r - darken + lighten))
        g = max(0, min(255, g - darken + lighten))
        b = max(0, min(255, b - darken + lighten))
        return f"#{r:02X}{g:02X}{b:02X}"

    @staticmethod
    def _muted_color(hex_color: str, opacity: float) -> str:
        """Blend *hex_color* toward white by *opacity* (0 = white, 1 = original).

        Returns a solid hex color approximating the visual result of the
        original color at the given opacity on a white background.
        """
        h = hex_color.lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        r = int(255 + (r - 255) * opacity)
        g = int(255 + (g - 255) * opacity)
        b = int(255 + (b - 255) * opacity)
        return f"#{r:02X}{g:02X}{b:02X}"

    def _set_slide_background(self, slide: Slide, color: str) -> None:
        """Set the solid background color of a slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(color)
