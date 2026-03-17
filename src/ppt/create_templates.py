"""PPT 模板生成脚本 -- 生成5套专业 .pptx 母版模板文件。

每个模板包含10种预制样板页（带装饰元素），渲染器在运行时复制对应页
并替换文本/图片内容。

用法::

    python -m src.ppt.create_templates
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Slide dimensions (16:9)
# ---------------------------------------------------------------------------
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

MARGIN_LEFT = Inches(0.9)
MARGIN_RIGHT = Inches(0.9)
MARGIN_TOP = Inches(0.7)
MARGIN_BOTTOM = Inches(0.7)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT

TEMPLATES_DIR = Path(__file__).parent / "templates"

# Layout index constants -- order of pages in the template file.
# Renderer uses these to locate the correct page to clone.
LAYOUT_TITLE = 0
LAYOUT_SECTION = 1
LAYOUT_TEXT_LEFT_IMG_RIGHT = 2
LAYOUT_IMG_LEFT_TEXT_RIGHT = 3
LAYOUT_FULL_IMAGE = 4
LAYOUT_THREE_COLUMNS = 5
LAYOUT_QUOTE = 6
LAYOUT_DATA_HIGHLIGHT = 7
LAYOUT_BLANK = 8
LAYOUT_CLOSING = 9
LAYOUT_COMPARISON = 10


# ---------------------------------------------------------------------------
# Theme definition dataclass
# ---------------------------------------------------------------------------


@dataclass
class TemplateTheme:
    name: str
    display_name: str
    primary: str
    accent: str
    secondary: str = ""
    bg: str = "#FFFFFF"
    text_color: str = "#2D3436"
    text_secondary: str = "#757575"
    cn_font: str = "Arial"  # cross-platform safe default
    en_font: str = "Arial"
    code_font: str = "Courier New"
    dark_mode: bool = False
    extra_colors: dict[str, str] = field(default_factory=dict)

    def __post_init__(self):
        if not self.secondary:
            self.secondary = self._muted(self.primary, 0.5)

    @staticmethod
    def _muted(hex_color: str, opacity: float) -> str:
        h = hex_color.lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        r = int(255 + (r - 255) * opacity)
        g = int(255 + (g - 255) * opacity)
        b = int(255 + (b - 255) * opacity)
        return f"#{r:02X}{g:02X}{b:02X}"


# ---------------------------------------------------------------------------
# 5 theme definitions
# ---------------------------------------------------------------------------

THEMES: list[TemplateTheme] = [
    TemplateTheme(
        name="modern",
        display_name="简约现代",
        primary="#2D3436",
        accent="#0984E3",
        secondary="#636E72",
        bg="#FFFFFF",
        text_color="#2D3436",
        text_secondary="#636E72",
    ),
    TemplateTheme(
        name="business",
        display_name="商务正式",
        primary="#1B2838",
        accent="#C0A062",
        secondary="#4A6278",
        bg="#FFFFFF",
        text_color="#1B2838",
        text_secondary="#5A6A7A",
    ),
    TemplateTheme(
        name="creative",
        display_name="创意活泼",
        primary="#FF6B6B",
        accent="#4ECDC4",
        secondary="#FFE66D",
        bg="#FFFFFF",
        text_color="#2C3E50",
        text_secondary="#7F8C8D",
    ),
    TemplateTheme(
        name="tech",
        display_name="科技极客",
        primary="#00D4FF",
        accent="#7C4DFF",
        secondary="#00E676",
        bg="#0D1117",
        text_color="#E6EDF3",
        text_secondary="#8B949E",
        dark_mode=True,
    ),
    TemplateTheme(
        name="education",
        display_name="教育清新",
        primary="#4A90E2",
        accent="#7ED321",
        secondary="#F5A623",
        bg="#F8F9FA",
        text_color="#2C3E50",
        text_secondary="#7F8C8D",
    ),
]


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    if len(h) != 6:
        h = "2D3436"
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_bg(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(color)


def _add_shape(
    slide,
    shape_type,
    left,
    top,
    width,
    height,
    *,
    fill_color: str | None = None,
    line_color: str | None = None,
    line_width: float | None = None,
):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = _hex_to_rgb(line_color)
        if line_width is not None:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_name: str = "Arial",
    font_size: int = 16,
    bold: bool = False,
    italic: bool = False,
    color: str = "#2D3436",
    alignment=PP_ALIGN.LEFT,
    line_spacing: float = 1.5,
    vertical_anchor=MSO_ANCHOR.TOP,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _hex_to_rgb(color)
    p.alignment = alignment

    # Line spacing
    pPr = p._p.get_or_add_pPr()
    lnSpc = pPr.makeelement(qn("a:lnSpc"), {})
    spcPct = lnSpc.makeelement(
        qn("a:spcPct"), {"val": str(int(line_spacing * 100000))}
    )
    lnSpc.append(spcPct)
    pPr.append(lnSpc)

    return txBox


def _add_gradient_rect(slide, left, top, width, height, color1: str, color2: str):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()

    spPr = shape._element.spPr
    for child in list(spPr):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("solidFill", "gradFill", "noFill", "pattFill"):
            spPr.remove(child)

    gradFill = spPr.makeelement(qn("a:gradFill"), {"flip": "none"})
    gsLst = gradFill.makeelement(qn("a:gsLst"), {})

    for pos, hex_color in [("0", color1), ("100000", color2)]:
        gs = gsLst.makeelement(qn("a:gs"), {"pos": pos})
        srgbClr = gs.makeelement(qn("a:srgbClr"), {"val": hex_color.lstrip("#")})
        gs.append(srgbClr)
        gsLst.append(gs)

    gradFill.append(gsLst)
    lin = gradFill.makeelement(qn("a:lin"), {"ang": "0", "scaled": "1"})
    gradFill.append(lin)
    spPr.append(gradFill)
    return shape


def _add_accent_line(slide, left, top, width, color: str):
    return _add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        left,
        top,
        width,
        Inches(0.03),
        fill_color=color,
    )


def _muted(hex_color: str, opacity: float) -> str:
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = int(255 + (r - 255) * opacity)
    g = int(255 + (g - 255) * opacity)
    b = int(255 + (b - 255) * opacity)
    return f"#{r:02X}{g:02X}{b:02X}"


def _add_overlay(slide, color: str = "#000000", opacity: float = 0.6):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(color)

    spPr = shape._element.spPr
    solidFill_elem = spPr.find(qn("a:solidFill"))
    if solidFill_elem is not None:
        srgbClr = solidFill_elem.find(qn("a:srgbClr"))
        if srgbClr is not None:
            alpha_val = str(int((1 - opacity) * 100000))
            alpha_elem = srgbClr.makeelement(qn("a:alpha"), {"val": alpha_val})
            srgbClr.append(alpha_elem)
    return shape


# ---------------------------------------------------------------------------
# Slide builders -- one per layout type
# ---------------------------------------------------------------------------


def _build_title_slide(prs: Presentation, t: TemplateTheme) -> None:
    """Title Slide (cover page)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, t.bg)

    # Accent gradient bar at bottom (15% height)
    bar_height = Inches(1.125)
    _add_gradient_rect(
        slide,
        Emu(0),
        SLIDE_HEIGHT - bar_height,
        SLIDE_WIDTH,
        bar_height,
        t.primary,
        t.accent,
    )

    # Title placeholder
    _add_textbox(
        slide,
        MARGIN_LEFT,
        Inches(2.2),
        CONTENT_WIDTH,
        Inches(1.2),
        "{title}",
        font_name=t.cn_font,
        font_size=42,
        bold=True,
        color=t.primary if not t.dark_mode else t.text_color,
        alignment=PP_ALIGN.CENTER,
    )

    # Subtitle placeholder
    _add_textbox(
        slide,
        MARGIN_LEFT,
        Inches(3.5),
        CONTENT_WIDTH,
        Inches(0.8),
        "{subtitle}",
        font_name=t.cn_font,
        font_size=22,
        color=t.text_color,
        alignment=PP_ALIGN.CENTER,
    )

    # Thin divider line
    line_width = Inches(2.0)
    line_left = (SLIDE_WIDTH - line_width) // 2
    _add_accent_line(slide, line_left, Inches(4.5), line_width, t.accent)

    # Metadata area
    _add_textbox(
        slide,
        MARGIN_LEFT,
        Inches(4.85),
        CONTENT_WIDTH,
        Inches(0.5),
        "{meta}",
        font_name=t.cn_font,
        font_size=13,
        color=t.text_secondary,
        alignment=PP_ALIGN.CENTER,
    )


def _build_section_slide(prs: Presentation, t: TemplateTheme) -> None:
    """Section Header slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Slightly different background
    if t.dark_mode:
        bg = "#161B22"
    else:
        h = t.bg.lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        r = max(0, r - 8)
        g = max(0, g - 8)
        b = max(0, b - 8)
        bg = f"#{r:02X}{g:02X}{b:02X}"
    _set_bg(slide, bg)

    center_y = Inches(3.0)
    line_width = Inches(2.0)
    line_left = (SLIDE_WIDTH - line_width) // 2

    # Top accent line
    _add_accent_line(slide, line_left, center_y, line_width, t.accent)

    # Section title
    _add_textbox(
        slide,
        MARGIN_LEFT,
        center_y + Inches(0.3),
        CONTENT_WIDTH,
        Inches(1.0),
        "{title}",
        font_name=t.cn_font,
        font_size=36,
        bold=True,
        color=t.primary if not t.dark_mode else t.text_color,
        alignment=PP_ALIGN.CENTER,
    )

    # Bottom accent line
    _add_accent_line(
        slide, line_left, center_y + Inches(1.5), line_width, t.accent
    )


def _build_content_with_image(
    prs: Presentation, t: TemplateTheme, *, text_on_left: bool
) -> None:
    """Content with Image (text + image split)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t.bg)

    text_width = Inches(6.3)
    image_width = Inches(4.6)
    gap = Inches(0.5)

    if text_on_left:
        text_left = MARGIN_LEFT
        img_left = MARGIN_LEFT + text_width + gap
    else:
        img_left = MARGIN_LEFT
        text_left = MARGIN_LEFT + image_width + gap

    # Title
    _add_textbox(
        slide,
        text_left,
        MARGIN_TOP,
        text_width,
        Inches(0.9),
        "{title}",
        font_name=t.cn_font,
        font_size=28,
        bold=True,
        color=t.primary if not t.dark_mode else t.text_color,
    )

    # Accent line under title
    _add_accent_line(slide, text_left, MARGIN_TOP + Inches(0.95), Inches(1.5), t.accent)

    # Content area placeholder
    _add_textbox(
        slide,
        text_left,
        MARGIN_TOP + Inches(1.4),
        text_width,
        Inches(4.5),
        "{content}",
        font_name=t.cn_font,
        font_size=18,
        color=t.text_color,
        line_spacing=1.6,
    )

    # Image placeholder (rounded rect)
    _add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        img_left,
        MARGIN_TOP + Inches(0.3),
        image_width,
        Inches(5.5),
        fill_color=_muted(t.accent, 0.12) if not t.dark_mode else _muted(t.accent, 0.2),
    )


def _build_full_image(prs: Presentation, t: TemplateTheme) -> None:
    """Full Image with overlay."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark gradient background as default
    _add_gradient_rect(
        slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, t.primary, "#1A1A2E"
    )

    # Dark overlay
    _add_overlay(slide, color="#000000", opacity=0.55)

    # Title
    _add_textbox(
        slide,
        Inches(1.5),
        Inches(3.5),
        Inches(10.333),
        Inches(1.2),
        "{title}",
        font_name=t.cn_font,
        font_size=38,
        bold=True,
        color="#FFFFFF",
        alignment=PP_ALIGN.CENTER,
    )

    # Body text
    _add_textbox(
        slide,
        Inches(2.0),
        Inches(4.9),
        Inches(9.333),
        Inches(1.0),
        "{content}",
        font_name=t.cn_font,
        font_size=20,
        color="#E0E0E0",
        alignment=PP_ALIGN.CENTER,
    )


def _build_three_columns(prs: Presentation, t: TemplateTheme) -> None:
    """Three Column layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t.bg)

    # Page title
    _add_textbox(
        slide,
        MARGIN_LEFT,
        MARGIN_TOP,
        CONTENT_WIDTH,
        Inches(0.9),
        "{title}",
        font_name=t.cn_font,
        font_size=28,
        bold=True,
        color=t.primary if not t.dark_mode else t.text_color,
        alignment=PP_ALIGN.CENTER,
    )

    # Divider
    _add_accent_line(
        slide,
        Inches(2.0),
        MARGIN_TOP + Inches(1.0),
        Inches(9.333),
        t.accent if not t.dark_mode else t.secondary,
    )

    # Three columns with decorative circles
    col_width = Inches(3.4)
    col_gap = Inches(0.37)
    col_start_left = MARGIN_LEFT + Inches(0.35)
    col_top = MARGIN_TOP + Inches(1.5)
    col_colors = [t.accent, t.secondary, t.primary]

    for i in range(3):
        col_left = col_start_left + Emu(int(col_width + col_gap) * i)

        # Decorative circle
        circle_size = Inches(0.55)
        circle_left = col_left + (col_width - circle_size) // 2
        _add_shape(
            slide,
            MSO_SHAPE.OVAL,
            circle_left,
            col_top,
            circle_size,
            circle_size,
            fill_color=col_colors[i],
        )

        # Column heading placeholder
        _add_textbox(
            slide,
            col_left,
            col_top + Inches(0.8),
            col_width,
            Inches(0.6),
            f"{{col{i+1}_title}}",
            font_name=t.cn_font,
            font_size=18,
            bold=True,
            color=t.primary if not t.dark_mode else t.text_color,
            alignment=PP_ALIGN.CENTER,
        )

        # Column description placeholder
        _add_textbox(
            slide,
            col_left,
            col_top + Inches(1.5),
            col_width,
            Inches(2.5),
            f"{{col{i+1}_desc}}",
            font_name=t.cn_font,
            font_size=15,
            color=t.text_color,
            alignment=PP_ALIGN.CENTER,
            line_spacing=1.5,
        )


def _build_quote(prs: Presentation, t: TemplateTheme) -> None:
    """Quote page."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t.bg)

    # Large decorative open-quote character
    _add_textbox(
        slide,
        Inches(1.5),
        Inches(1.0),
        Inches(2.0),
        Inches(1.8),
        "\u201C",
        font_name="Georgia",
        font_size=120,
        bold=True,
        color=_muted(t.accent, 0.35),
        alignment=PP_ALIGN.LEFT,
    )

    # Quote text placeholder
    _add_textbox(
        slide,
        Inches(2.0),
        Inches(2.5),
        Inches(9.333),
        Inches(2.5),
        "{quote}",
        font_name=t.cn_font,
        font_size=26,
        italic=True,
        color=t.text_color,
        alignment=PP_ALIGN.CENTER,
        line_spacing=1.7,
    )

    # Author placeholder
    _add_textbox(
        slide,
        Inches(2.0),
        Inches(5.3),
        Inches(9.333),
        Inches(0.5),
        "{author}",
        font_name=t.cn_font,
        font_size=16,
        color=t.text_secondary,
        alignment=PP_ALIGN.RIGHT,
    )


def _build_data_highlight(prs: Presentation, t: TemplateTheme) -> None:
    """Data Highlight page."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t.bg)

    # Page title
    _add_textbox(
        slide,
        MARGIN_LEFT,
        MARGIN_TOP,
        CONTENT_WIDTH,
        Inches(0.8),
        "{title}",
        font_name=t.cn_font,
        font_size=24,
        bold=True,
        color=t.primary if not t.dark_mode else t.text_color,
        alignment=PP_ALIGN.CENTER,
    )

    # Decorative circle behind number
    circle_size = Inches(3.0)
    circle_left = (SLIDE_WIDTH - circle_size) // 2
    _add_shape(
        slide,
        MSO_SHAPE.OVAL,
        circle_left,
        Inches(2.0),
        circle_size,
        circle_size,
        fill_color=_muted(t.accent, 0.12),
    )

    # Big data value placeholder
    _add_textbox(
        slide,
        MARGIN_LEFT,
        Inches(2.3),
        CONTENT_WIDTH,
        Inches(1.8),
        "{data_value}",
        font_name=t.cn_font,
        font_size=72,
        bold=True,
        color=t.accent,
        alignment=PP_ALIGN.CENTER,
    )

    # Data label placeholder
    _add_textbox(
        slide,
        MARGIN_LEFT,
        Inches(4.3),
        CONTENT_WIDTH,
        Inches(0.7),
        "{data_label}",
        font_name=t.cn_font,
        font_size=22,
        bold=True,
        color=t.primary if not t.dark_mode else t.text_color,
        alignment=PP_ALIGN.CENTER,
    )

    # Description placeholder
    _add_textbox(
        slide,
        Inches(2.5),
        Inches(5.2),
        Inches(8.333),
        Inches(0.8),
        "{description}",
        font_name=t.cn_font,
        font_size=16,
        color=t.text_secondary,
        alignment=PP_ALIGN.CENTER,
    )


def _build_blank(prs: Presentation, t: TemplateTheme) -> None:
    """Blank page (for timeline and other custom layouts)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t.bg)
    # Intentionally empty -- used for timeline and custom layouts


def _build_closing(prs: Presentation, t: TemplateTheme) -> None:
    """Closing page."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t.bg)

    # Bottom gradient bar
    bar_height = Inches(0.8)
    _add_gradient_rect(
        slide,
        Emu(0),
        SLIDE_HEIGHT - bar_height,
        SLIDE_WIDTH,
        bar_height,
        t.primary,
        t.accent,
    )

    # Main text placeholder
    _add_textbox(
        slide,
        MARGIN_LEFT,
        Inches(2.5),
        CONTENT_WIDTH,
        Inches(1.2),
        "{title}",
        font_name=t.cn_font,
        font_size=40,
        bold=True,
        color=t.primary if not t.dark_mode else t.text_color,
        alignment=PP_ALIGN.CENTER,
    )

    # Subtitle placeholder
    _add_textbox(
        slide,
        MARGIN_LEFT,
        Inches(3.8),
        CONTENT_WIDTH,
        Inches(0.7),
        "{subtitle}",
        font_name=t.cn_font,
        font_size=20,
        color=t.text_color,
        alignment=PP_ALIGN.CENTER,
    )

    # Contact info placeholder
    _add_textbox(
        slide,
        MARGIN_LEFT,
        Inches(4.8),
        CONTENT_WIDTH,
        Inches(0.5),
        "{contact}",
        font_name=t.cn_font,
        font_size=14,
        color=t.text_secondary,
        alignment=PP_ALIGN.CENTER,
    )


def _build_comparison(prs: Presentation, t: TemplateTheme) -> None:
    """Comparison page (two columns)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t.bg)

    # Title
    _add_textbox(
        slide,
        MARGIN_LEFT,
        MARGIN_TOP,
        CONTENT_WIDTH,
        Inches(0.9),
        "{title}",
        font_name=t.cn_font,
        font_size=28,
        bold=True,
        color=t.primary if not t.dark_mode else t.text_color,
        alignment=PP_ALIGN.CENTER,
    )

    # Two columns
    col_width = Inches(5.2)
    col_gap = Inches(1.0)
    col_top = MARGIN_TOP + Inches(1.3)
    col_height = Inches(5.0)
    left_col_left = MARGIN_LEFT + Inches(0.2)
    right_col_left = left_col_left + col_width + col_gap

    header_h = Inches(0.7)

    for idx, (col_left, col_color) in enumerate([
        (left_col_left, t.primary),
        (right_col_left, t.accent if not t.dark_mode else t.secondary),
    ]):
        # Header bar
        _add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            col_left,
            col_top,
            col_width,
            header_h,
            fill_color=col_color,
        )

        # Column header label
        label = "{left_title}" if idx == 0 else "{right_title}"
        _add_textbox(
            slide,
            col_left,
            col_top + Inches(0.05),
            col_width,
            Inches(0.5),
            label,
            font_name=t.cn_font,
            font_size=20,
            bold=True,
            color="#FFFFFF",
            alignment=PP_ALIGN.CENTER,
        )

        # Column content placeholder
        content_label = "{left_content}" if idx == 0 else "{right_content}"
        _add_textbox(
            slide,
            col_left + Inches(0.25),
            col_top + header_h + Inches(0.3),
            col_width - Inches(0.5),
            col_height - header_h - Inches(0.5),
            content_label,
            font_name=t.cn_font,
            font_size=17,
            color=t.text_color,
            line_spacing=1.6,
        )


# ---------------------------------------------------------------------------
# Main generation
# ---------------------------------------------------------------------------


def create_template(theme: TemplateTheme) -> Path:
    """Create a single template .pptx file and return its path."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Build slides in the canonical order (index matters for the renderer)
    _build_title_slide(prs, theme)           # 0: LAYOUT_TITLE
    _build_section_slide(prs, theme)         # 1: LAYOUT_SECTION
    _build_content_with_image(prs, theme, text_on_left=True)   # 2: LAYOUT_TEXT_LEFT_IMG_RIGHT
    _build_content_with_image(prs, theme, text_on_left=False)  # 3: LAYOUT_IMG_LEFT_TEXT_RIGHT
    _build_full_image(prs, theme)            # 4: LAYOUT_FULL_IMAGE
    _build_three_columns(prs, theme)         # 5: LAYOUT_THREE_COLUMNS
    _build_quote(prs, theme)                 # 6: LAYOUT_QUOTE
    _build_data_highlight(prs, theme)        # 7: LAYOUT_DATA_HIGHLIGHT
    _build_blank(prs, theme)                 # 8: LAYOUT_BLANK
    _build_closing(prs, theme)               # 9: LAYOUT_CLOSING
    _build_comparison(prs, theme)            # 10: LAYOUT_COMPARISON

    out_path = TEMPLATES_DIR / f"{theme.name}.pptx"
    prs.save(str(out_path))
    return out_path


def create_all_templates() -> list[Path]:
    """Generate all 5 template files and return their paths."""
    TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)
    paths: list[Path] = []
    for theme in THEMES:
        path = create_template(theme)
        paths.append(path)
        log.info("Created template: %s -> %s", theme.name, path)
    return paths


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO, format="%(message)s")
    paths = create_all_templates()
    for p in paths:
        print(f"  Created: {p}")
    print(f"\nAll {len(paths)} templates generated in {TEMPLATES_DIR}")
