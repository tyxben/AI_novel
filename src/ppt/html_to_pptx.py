"""HTML 幻灯片转 PPTX — 使用 Playwright 截图 + python-pptx 组装

将 M1 生成的 HTML 幻灯片文件转换为标准 .pptx 文件：
1. 解析 HTML 统计页数（<section class="slide">）
2. Playwright headless Chromium 逐页截图
3. BeautifulSoup 提取文本层（可选，用于备注/辅助功能）
4. python-pptx 将截图组装为全屏幻灯片，附加备注文本
"""

from __future__ import annotations

import logging
import re
import shutil
import tempfile
from pathlib import Path
from typing import Callable, Optional

from pptx import Presentation
from pptx.util import Inches

# ---------------------------------------------------------------------------
# Optional dependencies — lazy import with graceful fallback
# ---------------------------------------------------------------------------

try:
    from playwright.sync_api import sync_playwright  # noqa: F401

    _HAS_PLAYWRIGHT = True
except ImportError:
    _HAS_PLAYWRIGHT = False

try:
    from bs4 import BeautifulSoup  # noqa: F401

    _HAS_BS4 = True
except ImportError:
    _HAS_BS4 = False

log = logging.getLogger("ppt")

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# 16:9 standard widescreen dimensions (in inches)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Playwright viewport (pixels) — matches 16:9 aspect ratio
VIEWPORT_WIDTH = 1280
VIEWPORT_HEIGHT = 720

# Per-page screenshot timeout (ms)
PAGE_TIMEOUT_MS = 10_000

# Delay after showSlide() before screenshot (ms)
SLIDE_TRANSITION_DELAY_MS = 500


class HTMLToPPTXConverter:
    """HTML 幻灯片转 PPTX 转换器。

    Parameters
    ----------
    workspace : str
        工作目录路径，用于存放临时文件。
    extract_text : bool
        是否提取文本层写入幻灯片备注（需要 BeautifulSoup）。
    """

    def __init__(self, workspace: str = "workspace", extract_text: bool = True):
        if not _HAS_PLAYWRIGHT:
            raise ImportError(
                "Playwright is required for HTML-to-PPTX conversion. "
                "Install it with:\n"
                "  pip install playwright\n"
                "  playwright install chromium"
            )
        self.workspace = Path(workspace)
        self.extract_text = extract_text and _HAS_BS4

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def convert(
        self,
        html_path: str | Path,
        output_path: str | Path,
        progress_callback: Optional[Callable[[int, int, str], None]] = None,
    ) -> str:
        """将 HTML 幻灯片文件转换为 PPTX。

        Parameters
        ----------
        html_path : str | Path
            输入 HTML 文件路径。
        output_path : str | Path
            输出 PPTX 文件路径。
        progress_callback : callable, optional
            进度回调函数，签名: ``fn(page_num, total, message)``。

        Returns
        -------
        str
            生成的 PPTX 文件绝对路径。

        Raises
        ------
        FileNotFoundError
            如果 HTML 文件不存在。
        ValueError
            如果 HTML 中没有找到 ``<section class="slide">`` 元素。
        """
        html_path = Path(html_path)
        output_path = Path(output_path)

        if not html_path.exists():
            raise FileNotFoundError(f"HTML file not found: {html_path}")

        # Read HTML content
        html_content = html_path.read_text(encoding="utf-8")

        # Count total slides
        total_pages = self._count_slides(html_content)
        if total_pages == 0:
            raise ValueError(
                f"No slides found in {html_path}. "
                "Expected <section class=\"slide\"> elements."
            )

        log.info("Found %d slides in %s", total_pages, html_path.name)

        # Create temp directory for screenshots
        screenshots_dir = Path(tempfile.mkdtemp(prefix="ppt_screenshots_"))

        try:
            # Step 1: Capture screenshots
            screenshot_paths = self._capture_screenshots(
                html_path, screenshots_dir, total_pages, progress_callback
            )

            # Step 2: Extract text layers (optional)
            text_layers: list[str] = []
            if self.extract_text:
                text_layers = self._extract_text_layers(html_content)

            # Step 3: Assemble PPTX
            result_path = self._assemble_pptx(
                screenshot_paths, text_layers, output_path, progress_callback
            )

            log.info("PPTX saved to %s", result_path)
            return str(result_path)
        finally:
            # Step 4: Cleanup
            self._cleanup_screenshots(screenshots_dir)

    # ------------------------------------------------------------------
    # Internal methods
    # ------------------------------------------------------------------

    def _count_slides(self, html_content: str) -> int:
        """Count ``<section class="slide">`` elements in HTML content."""
        # Use regex to avoid requiring BS4 for counting
        pattern = r'<section\s+[^>]*class\s*=\s*["\'][^"\']*\bslide\b[^"\']*["\']'
        return len(re.findall(pattern, html_content, re.IGNORECASE))

    def _capture_screenshots(
        self,
        html_path: Path,
        output_dir: Path,
        total_pages: int,
        progress_callback: Optional[Callable[[int, int, str], None]],
    ) -> list[Path]:
        """Launch Playwright Chromium and capture per-slide screenshots.

        Parameters
        ----------
        html_path : Path
            HTML file to render.
        output_dir : Path
            Directory to save screenshot PNGs.
        total_pages : int
            Expected number of slides.
        progress_callback : callable, optional
            Progress callback.

        Returns
        -------
        list[Path]
            Ordered list of screenshot file paths.
        """
        from playwright.sync_api import sync_playwright

        screenshots: list[Path] = []
        file_url = html_path.resolve().as_uri()

        log.info("Launching Chromium for %d screenshots...", total_pages)

        with sync_playwright() as pw:
            browser = pw.chromium.launch(headless=True)
            try:
                page = browser.new_page(
                    viewport={"width": VIEWPORT_WIDTH, "height": VIEWPORT_HEIGHT}
                )

                # Load the HTML file
                page.goto(file_url, wait_until="networkidle", timeout=30_000)

                for i in range(total_pages):
                    page_num = i + 1

                    if progress_callback:
                        progress_callback(
                            page_num,
                            total_pages,
                            f"截图第 {page_num}/{total_pages} 页",
                        )

                    # Navigate to the target slide (0-based index)
                    page.evaluate(f"showSlide({i})")
                    page.wait_for_timeout(SLIDE_TRANSITION_DELAY_MS)

                    # Find the visible slide element and screenshot it
                    slide_selector = f'section.slide[data-page="{page_num}"]'
                    slide_element = page.query_selector(slide_selector)

                    # Fallback: try generic .slide selector if data-page not found
                    if slide_element is None:
                        slides = page.query_selector_all("section.slide")
                        if i < len(slides):
                            slide_element = slides[i]

                    # Final fallback: screenshot the whole viewport
                    screenshot_path = output_dir / f"slide_{page_num:03d}.png"

                    if slide_element is not None:
                        slide_element.screenshot(
                            path=str(screenshot_path),
                            timeout=PAGE_TIMEOUT_MS,
                        )
                    else:
                        log.warning(
                            "Slide element not found for page %d, "
                            "falling back to full-page screenshot",
                            page_num,
                        )
                        page.screenshot(path=str(screenshot_path), timeout=PAGE_TIMEOUT_MS)

                    screenshots.append(screenshot_path)
                    log.debug("Captured slide %d/%d", page_num, total_pages)
            finally:
                browser.close()

        log.info("All %d screenshots captured", len(screenshots))
        return screenshots

    def _extract_text_layers(self, html_content: str) -> list[str]:
        """Extract text content from each slide as markdown-like strings.

        Parameters
        ----------
        html_content : str
            Full HTML content containing ``<section class="slide">`` elements.

        Returns
        -------
        list[str]
            One markdown string per slide. Empty string for slides with no text.
        """
        if not _HAS_BS4:
            return []

        from bs4 import BeautifulSoup

        soup = BeautifulSoup(html_content, "html.parser")
        slides = soup.find_all("section", class_="slide")
        text_layers: list[str] = []

        for slide in slides:
            parts: list[str] = []

            # Extract h1 → "# title"
            h1 = slide.find("h1")
            if h1 and h1.get_text(strip=True):
                parts.append(f"# {h1.get_text(strip=True)}")

            # Extract h2 → "## heading"
            for h2 in slide.find_all("h2"):
                text = h2.get_text(strip=True)
                if text:
                    parts.append(f"## {text}")

            # Extract .subtitle → "## subtitle"
            for subtitle in slide.select(".subtitle"):
                text = subtitle.get_text(strip=True)
                if text:
                    parts.append(f"## {text}")

            # Extract li → "- bullet"
            for li in slide.find_all("li"):
                text = li.get_text(strip=True)
                if text:
                    parts.append(f"- {text}")

            # Extract p (not subtitle) → paragraph text
            for p in slide.find_all("p"):
                # Skip if this p is a subtitle (already captured)
                if p.get("class") and "subtitle" in p.get("class", []):
                    continue
                text = p.get_text(strip=True)
                if text:
                    parts.append(text)

            text_layers.append("\n".join(parts))

        return text_layers

    def _assemble_pptx(
        self,
        screenshot_paths: list[Path],
        text_layers: list[str],
        output_path: Path,
        progress_callback: Optional[Callable[[int, int, str], None]],
    ) -> Path:
        """Assemble screenshots into a PPTX file.

        Parameters
        ----------
        screenshot_paths : list[Path]
            Ordered list of slide screenshot files.
        text_layers : list[str]
            Optional text for each slide's notes. Can be shorter than
            screenshot_paths (missing entries treated as empty).
        output_path : Path
            Where to save the resulting .pptx file.
        progress_callback : callable, optional
            Progress callback.

        Returns
        -------
        Path
            The saved PPTX file path.
        """
        prs = Presentation()

        # Set 16:9 slide dimensions
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        # Use blank layout (index 6)
        blank_layout = prs.slide_layouts[6]

        total = len(screenshot_paths)
        for idx, img_path in enumerate(screenshot_paths):
            page_num = idx + 1

            if progress_callback:
                progress_callback(
                    page_num,
                    total,
                    f"组装第 {page_num}/{total} 页",
                )

            slide = prs.slides.add_slide(blank_layout)

            # Add full-screen image
            slide.shapes.add_picture(
                str(img_path),
                left=0,
                top=0,
                width=SLIDE_WIDTH,
                height=SLIDE_HEIGHT,
            )

            # Add notes text if available
            if idx < len(text_layers) and text_layers[idx]:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = text_layers[idx]

            log.debug("Assembled slide %d/%d", page_num, total)

        # Ensure output directory exists
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

        log.info("PPTX assembled: %d slides → %s", total, output_path)
        return output_path

    def _cleanup_screenshots(self, screenshots_dir: Path) -> None:
        """Remove temporary screenshot directory and all its contents."""
        if screenshots_dir.exists():
            shutil.rmtree(screenshots_dir, ignore_errors=True)
            log.debug("Cleaned up screenshots dir: %s", screenshots_dir)
